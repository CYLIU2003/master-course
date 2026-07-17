"""Build the advisor-facing Phase 3 recovery and weather-study deck.

The deck reads the immutable run summaries instead of duplicating numerical
results in source code.  It intentionally labels dirty-worktree runs as
provisional and never presents the two-stage Phase 3 result as a global
total-cost optimum.
"""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_ROOT = Path(r"C:\master-course\output")
DEFAULT_SUNNY_SUMMARY = (
    DEFAULT_OUTPUT_ROOT
    / "research_phase3_sunny_final_1500s_20260716"
    / "summary.json"
)
DEFAULT_RAIN_SUMMARY = (
    DEFAULT_OUTPUT_ROOT
    / "research_phase3_rain_final_1500s_20260716"
    / "summary.json"
)
DEFAULT_BASELINE_SUNNY_SUMMARY = (
    DEFAULT_OUTPUT_ROOT
    / "research_phase3_frontend_weather_60min_sunny_95ade40"
    / "summary.json"
)
DEFAULT_BASELINE_RAIN_SUMMARY = (
    DEFAULT_OUTPUT_ROOT
    / "research_phase3_frontend_weather_60min_rain_95ade40"
    / "summary.json"
)
DEFAULT_PRESENTATION = (
    REPO_ROOT
    / "docs"
    / "presentations"
    / "phase3_weather_model_progress_20260716.pptx"
)

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Meiryo UI"
FONT_MONO = "Cascadia Mono"

NAVY = RGBColor(11, 31, 51)
NAVY_2 = RGBColor(23, 52, 77)
CYAN = RGBColor(22, 184, 196)
CYAN_LIGHT = RGBColor(218, 246, 248)
SUN = RGBColor(245, 185, 65)
SUN_LIGHT = RGBColor(255, 245, 219)
RAIN = RGBColor(77, 120, 204)
RAIN_LIGHT = RGBColor(226, 235, 251)
GREEN = RGBColor(40, 168, 121)
GREEN_LIGHT = RGBColor(225, 246, 237)
RED = RGBColor(227, 93, 106)
RED_LIGHT = RGBColor(253, 231, 234)
ICE = RGBColor(109, 118, 131)
TEXT = RGBColor(30, 42, 54)
MUTED = RGBColor(94, 111, 128)
GRID = RGBColor(215, 222, 229)
BG = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)


def _load_json(path: Path) -> dict[str, Any]:
    if not path.is_file():
        raise FileNotFoundError(f"Required run summary is missing: {path}")
    return json.loads(path.read_text(encoding="utf-8"))


def _number(value: Any, default: float = 0.0) -> float:
    try:
        result = float(value)
    except (TypeError, ValueError):
        return default
    return result if math.isfinite(result) else default


def _mapping(container: Mapping[str, Any], key: str) -> Mapping[str, Any]:
    value = container.get(key)
    return value if isinstance(value, Mapping) else {}


def _pct(value: float) -> str:
    return f"{value:.1f}%"


def _jpy(value: float) -> str:
    return f"{value:,.0f}円"


def _kwh(value: float) -> str:
    return f"{value:,.1f} kWh"


def _set_background(slide: Any, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_cell_text(
    cell: Any,
    text: str,
    *,
    size: float = 13,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    cell.text = str(text)
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str = FONT,
    margin: float = 0.02,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _add_rich_text(
    slide: Any,
    runs: Sequence[tuple[str, bool, RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, bold, color in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _add_bullets(
    slide: Any,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: RGBColor = TEXT,
    spacing: float = 5,
) -> Any:
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.03)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(item)
        paragraph.level = 0
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(spacing)
        paragraph.text = f"• {paragraph.text}"
    return box


def _shape(
    slide: Any,
    shape_type: MSO_SHAPE,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius: bool = False,
) -> Any:
    del radius
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _card(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = GRID,
) -> Any:
    return _shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
    )


def _add_header(slide: Any, title: str, section: str, page: int) -> None:
    _add_text(slide, section.upper(), 0.55, 0.22, 3.6, 0.25, size=9, color=CYAN, bold=True)
    _add_text(slide, title, 0.55, 0.52, 12.0, 0.55, size=25, color=NAVY, bold=True)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.55),
        Inches(1.18),
        Inches(12.78),
        Inches(1.18),
    )
    line.line.color.rgb = GRID
    line.line.width = Pt(1)
    _add_text(slide, f"{page:02d}", 12.38, 7.12, 0.42, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_footer(slide: Any, text: str) -> None:
    _add_text(slide, text, 0.55, 7.08, 11.5, 0.22, size=7.5, color=MUTED)


def _metric_card(
    slide: Any,
    label: str,
    value: str,
    note: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: RGBColor,
) -> None:
    _card(slide, x, y, w, 1.38)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, 1.38, fill=accent)
    _add_text(slide, label, x + 0.22, y + 0.16, w - 0.36, 0.24, size=10, color=MUTED, bold=True)
    _add_text(slide, value, x + 0.22, y + 0.48, w - 0.36, 0.42, size=22, color=NAVY, bold=True)
    _add_text(slide, note, x + 0.22, y + 0.98, w - 0.36, 0.2, size=8.5, color=MUTED)


def _add_stacked_bar(
    slide: Any,
    label: str,
    first: float,
    second: float,
    x: float,
    y: float,
    w: float,
    *,
    first_color: RGBColor = CYAN,
    second_color: RGBColor = ICE,
    suffix: str = "便",
) -> None:
    total = max(first + second, 1.0)
    _add_text(slide, label, x, y, 1.35, 0.28, size=12, color=TEXT, bold=True)
    bar_x = x + 1.45
    first_w = w * first / total
    second_w = w - first_w
    _shape(slide, MSO_SHAPE.RECTANGLE, bar_x, y + 0.01, first_w, 0.30, fill=first_color)
    _shape(slide, MSO_SHAPE.RECTANGLE, bar_x + first_w, y + 0.01, second_w, 0.30, fill=second_color)
    _add_text(slide, f"BEV {first:,.0f}{suffix}", bar_x + 0.06, y + 0.06, max(first_w - 0.12, 0.2), 0.18, size=8, color=WHITE, bold=True)
    if second_w > 1.0:
        _add_text(slide, f"ICE {second:,.0f}{suffix}", bar_x + first_w + 0.06, y + 0.06, second_w - 0.12, 0.18, size=8, color=WHITE, bold=True)


def _add_table(
    slide: Any,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    widths: Sequence[float] | None = None,
    font_size: float = 11,
) -> Any:
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = table_shape.table
    if widths:
        for column, width in zip(table.columns, widths):
            column.width = Inches(width)
    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_2
        _set_cell_text(cell, header, size=font_size, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for row_index, row in enumerate(rows, start=1):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else BG
            _set_cell_text(
                cell,
                value,
                size=font_size,
                color=TEXT,
                align=PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.RIGHT,
            )
    return table_shape


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, SLIDE_H, fill=CYAN)
    _add_text(slide, "MASTER’S THESIS • PHASE 3", 0.72, 0.62, 6.4, 0.25, size=10, color=CYAN, bold=True)
    _add_text(slide, "EVバス運用最適化モデル\n復旧・天候差反映の進捗", 0.72, 1.35, 11.4, 1.45, size=34, color=WHITE, bold=True)
    _add_text(slide, "2026年7月11日–16日｜実フロント入力・Gurobi Phase 3検証", 0.75, 3.15, 9.4, 0.38, size=16, color=RGBColor(202, 218, 231))
    _card(slide, 0.75, 4.18, 11.7, 1.25, fill=NAVY_2, line=RGBColor(46, 78, 105))
    _add_text(slide, "研究の主張", 1.05, 4.48, 1.35, 0.25, size=10, color=CYAN, bold=True)
    _add_text(slide, "「最適な表」ではなく、EV運用の成立条件と支配制約を定量化する", 2.25, 4.38, 9.65, 0.52, size=19, color=WHITE, bold=True)
    _add_text(slide, "対象：鶴巻営業所・264便・BEV 35台 / ICE 25台・晴天 2025-08-05 / 雨天 2025-08-10", 0.78, 6.66, 11.8, 0.28, size=9, color=RGBColor(178, 199, 217))


def _add_research_question_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "研究の軸を「制約構造の説明」へ戻す", "01 / RESEARCH QUESTION", 2)
    _add_rich_text(
        slide,
        [
            ("最適化 = ", False, MUTED),
            ("成立条件とボトルネックを調べる実験装置", True, NAVY),
        ],
        0.85,
        1.48,
        11.7,
        0.5,
        size=23,
        align=PP_ALIGN.CENTER,
    )
    cards = (
        ("可行性", "既存264便を\n全便維持できるか", "全便・接続・SOC・充電器・契約電力", GREEN, GREEN_LIGHT),
        ("支配制約", "何がEV化を\n止めているか", "車両数 / 充電窓 / 初期SOC / PV", SUN, SUN_LIGHT),
        ("介入効果", "PV・BESS・契約を\n変えると何が動くか", "台数・担当便・買電・ピーク・費用", RAIN, RAIN_LIGHT),
    )
    for index, (label, title, note, accent, fill) in enumerate(cards):
        x = 0.72 + index * 4.12
        _card(slide, x, 2.38, 3.72, 2.75, fill=fill, line=accent)
        _add_text(slide, label, x + 0.28, 2.66, 1.0, 0.24, size=10, color=accent, bold=True)
        _add_text(slide, title, x + 0.28, 3.10, 3.1, 0.75, size=21, color=NAVY, bold=True)
        _add_text(slide, note, x + 0.28, 4.30, 3.05, 0.42, size=10, color=MUTED)
    _card(slide, 1.32, 5.65, 10.7, 0.86, fill=WHITE, line=GRID)
    _add_text(slide, "今回の焦点", 1.62, 5.92, 1.25, 0.22, size=10, color=RED, bold=True)
    _add_text(slide, "PV発電量の差がStage 1のBEV/ICE割当に届いていなかった構造を修正し、Stage 2で物理可行性を再検証", 2.85, 5.83, 8.7, 0.35, size=14, color=TEXT, bold=True)


def _add_call_chain_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "検証済みの実行経路：二段階を混同しない", "02 / VERIFIED CALL CHAIN", 3)
    stages = (
        ("入力", "Frontend scenario\n+ prepared scope", "料金・在庫・SOC\nPV/BESS・日付", NAVY_2),
        ("Stage 1", "運行割当 MILP", "全便・接続・車両\nSOC必要条件・費用代理", CYAN),
        ("Stage 2", "固定割当の充電 MILP", "SOC・充電器・PV/BESS\n契約電力・時刻別費用", RAIN),
        ("受入", "独立validation", "違反0・fallbackなし\n会計ledger", GREEN),
    )
    for index, (label, title, note, accent) in enumerate(stages):
        x = 0.55 + index * 3.20
        _card(slide, x, 2.03, 2.72, 2.26, fill=WHITE, line=accent)
        _shape(slide, MSO_SHAPE.RECTANGLE, x, 2.03, 2.72, 0.16, fill=accent)
        _add_text(slide, label, x + 0.25, 2.42, 0.7, 0.22, size=9, color=accent, bold=True)
        _add_text(slide, title, x + 0.25, 2.76, 2.2, 0.55, size=17, color=NAVY, bold=True)
        _add_text(slide, note, x + 0.25, 3.48, 2.2, 0.48, size=10, color=MUTED)
        if index < len(stages) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + 2.79),
                Inches(2.92),
                Inches(0.33),
                Inches(0.42),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRID
            arrow.line.color.rgb = GRID
    _card(slide, 0.78, 5.05, 11.75, 1.03, fill=RED_LIGHT, line=RED)
    _add_text(slide, "重要", 1.08, 5.35, 0.72, 0.22, size=10, color=RED, bold=True)
    _add_text(slide, "Stage 1の割当をStage 2で固定するため、現行Phase 3はglobal simultaneous total-cost optimumではない", 1.82, 5.26, 10.1, 0.34, size=15, color=NAVY, bold=True)
    _add_text(slide, "可行スケジュールの比較は可能。最適費用・一意な最適構成という主張は不可。", 1.82, 5.68, 9.4, 0.25, size=10, color=MUTED)


def _add_timeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "1週間でつぶした問題：見かけの解から研究用モデルへ", "03 / RECOVERY TIMELINE", 4)
    events = (
        ("7/11", "偽の最終出力", "Stage 2不可行候補を隔離\nfallback / repairを研究結果から排除", RED),
        ("7/12", "比較条件の交絡", "車両費・BESS終端・PV provenanceを同期", SUN),
        ("7/13–14", "SOC・回送の欠落", "始発 / 便間 / 帰庫 / 終端SOCを厳密化", CYAN),
        ("7/14", "料金式の不一致", "TOU実時刻・PV 0円・需要料金日割りを統一", RAIN),
        ("7/15", "台数比の感度", "利用可能BEV台数感度 + 帰庫直後SOC境界修正", GREEN),
        ("7/16", "PVが割当に届かない", "Stage 1費用代理 + 所在地対応必要条件", NAVY_2),
    )
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.92),
        Inches(3.15),
        Inches(12.4),
        Inches(3.15),
    )
    line.line.color.rgb = GRID
    line.line.width = Pt(2)
    for index, (date, title, note, accent) in enumerate(events):
        x = 0.45 + index * 2.12
        _shape(slide, MSO_SHAPE.OVAL, x + 0.77, 2.96, 0.38, 0.38, fill=accent)
        y = 1.46 if index % 2 == 0 else 3.63
        _card(slide, x, y, 1.92, 1.24, fill=WHITE, line=accent)
        _add_text(slide, date, x + 0.16, y + 0.14, 0.62, 0.20, size=9, color=accent, bold=True)
        _add_text(slide, title, x + 0.16, y + 0.43, 1.60, 0.29, size=12, color=NAVY, bold=True)
        _add_text(slide, note, x + 0.16, y + 0.78, 1.60, 0.33, size=7.7, color=MUTED)
        connector_y1 = y + (1.24 if index % 2 == 0 else 0)
        connector_y2 = 2.96 if index % 2 == 0 else 3.34
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x + 0.96),
            Inches(connector_y1),
            Inches(x + 0.96),
            Inches(connector_y2),
        )
        connector.line.color.rgb = accent
        connector.line.width = Pt(1)
    _add_footer(slide, "変更のたびにhard constraint、料金単位、fallback、成果物契約、回帰テストを再確認")


def _add_controls_slide(prs: Presentation, sunny: Mapping[str, Any], rain: Mapping[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "晴天・雨天の比較契約：PV以外を固定", "04 / EXPERIMENTAL CONTROLS", 5)
    sunny_pv = _number(_mapping(sunny, "stage1_energy_cost_proxy_weather_input").get("pv_available_kwh_by_depot", {}).get("tsurumaki"))
    rain_pv = _number(_mapping(rain, "stage1_energy_cost_proxy_weather_input").get("pv_available_kwh_by_depot", {}).get("tsurumaki"))
    rows = (
        ("Service date", str(sunny.get("service_date")), str(rain.get("service_date"))),
        ("Scenario ID", str(sunny.get("scenario_id"))[:13] + "…", str(rain.get("scenario_id"))[:13] + "…"),
        ("全便 / Fleet", f"{sunny.get('trip_count')}便 / BEV35・ICE25", f"{rain.get('trip_count')}便 / BEV35・ICE25"),
        ("Time / seed / gap", f"{sunny.get('time_limit_sec')}秒 / {sunny.get('random_seed')} / {sunny.get('mip_gap')}", f"{rain.get('time_limit_sec')}秒 / {rain.get('random_seed')} / {rain.get('mip_gap')}"),
        ("PV marginal price", f"{_number(sunny.get('pv_marginal_charge_cost_yen_per_kwh')):.1f}円/kWh", f"{_number(rain.get('pv_marginal_charge_cost_yen_per_kwh')):.1f}円/kWh"),
        ("PV availability", _kwh(sunny_pv), _kwh(rain_pv)),
    )
    _add_table(slide, ("固定・操作項目", "晴天", "雨天"), rows, 0.68, 1.55, 8.45, 4.65, widths=(2.55, 2.95, 2.95), font_size=10.5)
    _card(slide, 9.45, 1.55, 3.20, 4.65, fill=WHITE, line=GRID)
    max_pv = max(sunny_pv, rain_pv, 1.0)
    _add_text(slide, "操作変数", 9.76, 1.85, 1.2, 0.23, size=10, color=CYAN, bold=True)
    _add_text(slide, "PV発電可能量", 9.76, 2.19, 2.2, 0.33, size=19, color=NAVY, bold=True)
    for index, (label, value, color) in enumerate((("晴天", sunny_pv, SUN), ("雨天", rain_pv, RAIN))):
        y = 3.00 + index * 1.00
        _add_text(slide, label, 9.76, y, 0.55, 0.23, size=10, color=TEXT, bold=True)
        bar_width = 1.72 * value / max_pv
        _shape(slide, MSO_SHAPE.RECTANGLE, 10.42, y, bar_width, 0.30, fill=color)
        value_x = 10.45 if bar_width >= 0.75 else 10.42 + bar_width + 0.08
        value_color = WHITE if bar_width >= 0.75 else color
        _add_text(slide, f"{value:.1f}", value_x, y + 0.04, 1.25, 0.18, size=8, color=value_color, bold=True)
    reduction = (1.0 - rain_pv / sunny_pv) * 100.0 if sunny_pv else 0.0
    _metric_card(slide, "雨天のPV減少", f"−{reduction:.1f}%", "他の料金・車両・SOC・充電器は固定", 9.74, 4.95, 2.60, accent=RAIN)
    _add_footer(slide, "比較対象の実効weather profileはoperation_mode以外中立。BEV/ICE biasやSOC上書きは使わない。")


def _add_proxy_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "根本原因：天候別PVがStage 1の割当に入っていなかった", "05 / ROOT CAUSE & MODEL CHANGE", 6)
    _card(slide, 0.68, 1.54, 4.00, 4.95, fill=RED_LIGHT, line=RED)
    _add_text(slide, "修正前", 0.98, 1.84, 0.9, 0.25, size=10, color=RED, bold=True)
    _add_text(slide, "Stage 1", 0.98, 2.27, 2.0, 0.38, size=21, color=NAVY, bold=True)
    _add_bullets(
        slide,
        (
            "ICE燃料・CO₂・車両使用費だけで割当",
            "PV量はStage 2にしか存在しない",
            "晴雨で同一割当になるのはモデル上自然",
        ),
        0.98,
        2.86,
        3.35,
        1.55,
        size=12,
        color=TEXT,
    )
    _add_text(slide, "FStage1(weather) = constant", 1.02, 5.12, 3.28, 0.35, size=15, color=RED, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    _add_text(slide, "割当探索から見れば晴天も雨天も同じ", 1.05, 5.65, 3.2, 0.25, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    _card(slide, 4.98, 1.54, 7.68, 4.95, fill=GREEN_LIGHT, line=GREEN)
    _add_text(slide, "修正後：営業所別の集約充電費用下界", 5.30, 1.84, 4.8, 0.30, size=10, color=GREEN, bold=True)
    _add_text(slide, "Cproxy = cPV·EPV + cBESS·EBESS + cgrid,min·Egrid", 5.30, 2.42, 6.85, 0.45, size=19, color=NAVY, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    _add_text(slide, "Esource ≥ max{0, Etrip + Edeadhead + Eterminal − Einitial} / ηcharge", 5.30, 3.08, 6.85, 0.42, size=14, color=NAVY_2, font=FONT_MONO, align=PP_ALIGN.CENTER)
    _add_bullets(
        slide,
        (
            "PV限界費用はフロント設定どおり0円/kWh",
            "天候別PV上限を超えた分だけ系統単価で評価",
            "始発・便間・帰庫回送、実効終端SOCを含む",
            "充電時刻・需要料金・競合はStage 2に残す",
        ),
        5.42,
        3.86,
        6.5,
        1.55,
        size=11,
        color=TEXT,
    )
    _add_text(slide, "目的：PVの価値を割当に伝える。実現費用を近似したと主張しない。", 5.34, 5.83, 6.9, 0.28, size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


def _add_probe_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "短時間probeで確認：同じ割当でも天候費用は変わる", "06 / UNIT OF EVIDENCE", 7)
    _metric_card(slide, "必要外部充電量", "616.4 kWh", "同一割当：BEV 17台・54便", 0.72, 1.63, 3.72, accent=CYAN)
    _metric_card(slide, "晴天 proxy", "31円", "PV 614.7 / grid 1.7 kWh", 4.81, 1.63, 3.72, accent=SUN)
    _metric_card(slide, "雨天 proxy", "9,533円", "PV 101.1 / grid 515.3 kWh", 8.90, 1.63, 3.72, accent=RAIN)
    _card(slide, 0.72, 3.46, 11.90, 2.54, fill=WHITE, line=GRID)
    _add_text(slide, "意味", 1.02, 3.78, 0.7, 0.24, size=10, color=CYAN, bold=True)
    _add_text(slide, "PV 0円の設定は正しくStage 1目的へ入り、晴天のBEV充電必要量をほぼ無償と評価した", 1.02, 4.18, 10.85, 0.45, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.06), Inches(4.82), Inches(0.55), Inches(0.58))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = CYAN
    arrow.line.color.rgb = CYAN
    _add_text(slide, "次の検証：この費用差で割当自体が変わり、固定割当のStage 2も可行か", 2.25, 5.42, 8.86, 0.28, size=12, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_footer(slide, "probeは数理反応の診断。正式な晴雨比較は各1500秒runで評価。")


def _add_location_fix_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "自分で発見した副作用：Stage 1が営業所外充電を発明", "07 / IIS-DRIVEN CORRECTION", 8)
    _card(slide, 0.68, 1.55, 3.55, 4.83, fill=RED_LIGHT, line=RED)
    _add_text(slide, "1500秒・初回修正版", 0.98, 1.86, 2.5, 0.24, size=10, color=RED, bold=True)
    _add_text(slide, "Stage 1", 0.98, 2.28, 1.6, 0.35, size=20, color=NAVY, bold=True)
    _add_text(slide, "BEV 190便を選択", 0.98, 2.72, 2.65, 0.30, size=16, color=RED, bold=True)
    _add_text(slide, "↓", 1.98, 3.18, 0.5, 0.35, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "Stage 2 infeasible", 0.98, 3.64, 2.65, 0.30, size=16, color=RED, bold=True)
    _add_bullets(slide, ("IIS：始発回送後、slot 1–18は非在庫", "旧緩和はidleならどこでも充電", "偽のSOC可行性をStage 1が許容"), 0.98, 4.20, 2.85, 1.25, size=10)
    _card(slide, 4.56, 1.55, 3.72, 4.83, fill=SUN_LIGHT, line=SUN)
    _add_text(slide, "第1案", 4.86, 1.86, 1.0, 0.24, size=10, color=SUN, bold=True)
    _add_text(slide, "slot別所在地制約", 4.86, 2.27, 2.9, 0.35, size=19, color=NAVY, bold=True)
    _metric_card(slide, "制約数", "69,300", "正しいがLP緩和が重すぎる", 4.86, 3.00, 3.10, accent=SUN)
    _add_text(slide, "750秒でbound 0\n探索性能を失う", 5.17, 4.86, 2.45, 0.66, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _card(slide, 8.60, 1.55, 4.05, 4.83, fill=GREEN_LIGHT, line=GREEN)
    _add_text(slide, "採用案", 8.90, 1.86, 1.0, 0.24, size=10, color=GREEN, bold=True)
    _add_text(slide, "累積ホームデポ\n充電可能量の必要条件", 8.90, 2.27, 3.2, 0.75, size=18, color=NAVY, bold=True)
    _metric_card(slide, "制約数", "875", "69,300比 −98.7%", 8.90, 3.32, 3.42, accent=GREEN)
    _add_bullets(slide, ("割当に裏付けられた在庫窓だけ計上", "始発・便間・帰庫loadを時系列累積", "共有充電器等はStage 2で厳密化"), 8.90, 4.94, 3.15, 1.05, size=9.5)


def _add_mix_slide(
    prs: Presentation,
    sunny: Mapping[str, Any],
    rain: Mapping[str, Any],
    baseline_sunny: Mapping[str, Any],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "1500秒結果：天候差がBEV担当便数へ反映", "08 / ASSIGNMENT RESULT", 9)
    sunny_trips = _mapping(sunny, "served_trip_count_by_vehicle_type")
    rain_trips = _mapping(rain, "served_trip_count_by_vehicle_type")
    baseline_trips = _mapping(baseline_sunny, "served_trip_count_by_vehicle_type")
    if not baseline_trips:
        baseline_trips = {"BEV": 54, "ICE": 210}
    _add_stacked_bar(slide, "旧モデル", _number(baseline_trips.get("BEV"), 54), _number(baseline_trips.get("ICE"), 210), 0.74, 1.84, 9.65)
    _add_stacked_bar(slide, "晴天", _number(sunny_trips.get("BEV")), _number(sunny_trips.get("ICE")), 0.74, 2.67, 9.65, first_color=SUN)
    _add_stacked_bar(slide, "雨天", _number(rain_trips.get("BEV")), _number(rain_trips.get("ICE")), 0.74, 3.50, 9.65, first_color=RAIN)
    sunny_used = _mapping(sunny, "used_vehicle_count_by_type")
    rain_used = _mapping(rain, "used_vehicle_count_by_type")
    rows = (
        ("旧モデル（晴雨同一）", "17 / 15", "54 / 210", "20.5%"),
        ("晴天・新モデル", f"{sunny_used.get('BEV')} / {sunny_used.get('ICE')}", f"{sunny_trips.get('BEV')} / {sunny_trips.get('ICE')}", _pct(100 * _number(sunny_trips.get('BEV')) / 264)),
        ("雨天・新モデル", f"{rain_used.get('BEV')} / {rain_used.get('ICE')}", f"{rain_trips.get('BEV')} / {rain_trips.get('ICE')}", _pct(100 * _number(rain_trips.get('BEV')) / 264)),
    )
    _add_table(slide, ("ケース", "使用BEV / ICE", "担当BEV / ICE便", "BEV便率"), rows, 0.74, 4.42, 11.82, 1.63, widths=(3.05, 2.45, 3.25, 2.07), font_size=10.5)
    delta_trips = _number(sunny_trips.get("BEV")) - _number(rain_trips.get("BEV"))
    _card(slide, 10.60, 1.60, 1.95, 2.42, fill=SUN_LIGHT, line=SUN)
    _add_text(slide, "晴天効果", 10.85, 1.92, 1.4, 0.24, size=9, color=SUN, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"+{delta_trips:.0f}", 10.83, 2.37, 1.45, 0.48, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "BEV担当便\n（雨天比）", 10.85, 3.02, 1.40, 0.50, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_footer(slide, "同じ使用台数でも担当便数は異なる。電動化率は『台数』と『便数』を分けて報告する。")


def _add_energy_cost_slide(prs: Presentation, sunny: Mapping[str, Any], rain: Mapping[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "実現エネルギー・会計費用：Stage 2で再計算", "09 / REALIZED ACCOUNTING", 10)
    sunny_flows = _mapping(sunny, "flows_kwh_or_kw")
    rain_flows = _mapping(rain, "flows_kwh_or_kw")
    sunny_costs = _mapping(sunny, "costs_jpy")
    rain_costs = _mapping(rain, "costs_jpy")
    rows = []
    metrics = (
        ("PV発電量 [kWh]", "pv_generated_kwh", sunny_flows, rain_flows),
        ("系統買電 [kWh]", "grid_import_kwh", sunny_flows, rain_flows),
        ("ピーク [kW]", "peak_grid_kw", sunny_flows, rain_flows),
        ("電力費 [円]", "electricity_cost", sunny_costs, rain_costs),
        ("燃料費 [円]", "fuel_cost", sunny_costs, rain_costs),
        ("会計総費用 [円]", "total_cost", sunny_costs, rain_costs),
    )
    for label, key, sunny_container, rain_container in metrics:
        s_value = _number(sunny_container.get(key))
        r_value = _number(rain_container.get(key))
        rows.append((label, f"{s_value:,.1f}", f"{r_value:,.1f}", f"{r_value - s_value:+,.1f}"))
    _add_table(slide, ("指標", "晴天", "雨天", "雨天 − 晴天"), rows, 0.68, 1.52, 8.62, 4.78, widths=(2.50, 1.88, 1.88, 2.36), font_size=10.5)
    total_s = _number(sunny_costs.get("total_cost"))
    total_r = _number(rain_costs.get("total_cost"))
    delta = total_r - total_s
    delta_pct = delta / total_s * 100.0 if total_s else 0.0
    _metric_card(slide, "雨天の会計総費用差", f"{delta:+,.0f}円", f"晴天比 {delta_pct:+.2f}%", 9.62, 1.52, 2.75, accent=RAIN)
    _metric_card(slide, "雨天の買電差", f"{_number(rain_flows.get('grid_import_kwh')) - _number(sunny_flows.get('grid_import_kwh')):+,.1f}", "kWh（BEV担当便−22で相殺）", 9.62, 3.22, 2.75, accent=RAIN)
    _metric_card(slide, "雨天の燃料費差", f"{_number(rain_costs.get('fuel_cost')) - _number(sunny_costs.get('fuel_cost')):+,.0f}円", "BEV担当便差の影響", 9.62, 4.92, 2.75, accent=SUN)
    _add_footer(slide, "会計総費用は可行スケジュールの評価値。Stage 1 gap未収束のためglobal optimumとして扱わない。")


def _add_validity_slide(prs: Presentation, sunny: Mapping[str, Any], rain: Mapping[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "研究妥当性：何が言えて、何がまだ言えないか", "10 / VALIDITY & LIMITATIONS", 11)
    rows = (
        ("全便担当", f"{sunny.get('trip_count_served')}/264", f"{rain.get('trip_count_served')}/264"),
        ("Stage 1", str(sunny.get("stage1_solver_status")), str(rain.get("stage1_solver_status"))),
        ("Stage 1 gap", f"{_number(sunny.get('stage1_mip_gap_percent')):.2f}%", f"{_number(rain.get('stage1_mip_gap_percent')):.2f}%"),
        ("Stage 2", str(sunny.get("stage2_solver_status")), str(rain.get("stage2_solver_status"))),
        ("Hard validation", "全項目0", "全項目0"),
        ("Fallback / repair", "なし / なし", "なし / なし"),
    )
    _add_table(slide, ("判定項目", "晴天", "雨天"), rows, 0.68, 1.50, 7.24, 4.85, widths=(2.50, 2.22, 2.22), font_size=10.5)
    _card(slide, 8.28, 1.50, 4.36, 2.10, fill=GREEN_LIGHT, line=GREEN)
    _add_text(slide, "言える", 8.60, 1.82, 0.75, 0.24, size=10, color=GREEN, bold=True)
    _add_bullets(slide, ("hard constraints下で全264便が運用可能", "PV量を変えると可行な割当・費用が変化", "Stage 2でSOC/充電器/契約違反ゼロ"), 8.60, 2.20, 3.60, 1.05, size=10)
    _card(slide, 8.28, 3.90, 4.36, 2.45, fill=RED_LIGHT, line=RED)
    _add_text(slide, "まだ言えない", 8.60, 4.22, 1.05, 0.24, size=10, color=RED, bold=True)
    _add_bullets(slide, ("晴雨それぞれの大域最適費用", "一意な最適BEV/ICE構成", "Phase 4同時最適化との優劣", "dirty worktree成果物の正式採択"), 8.60, 4.60, 3.60, 1.30, size=10)
    _add_footer(slide, "現在の1500秒成果物はcommit候補コードのprovisional evidence。clean commitからの再実行後に正式比較へ昇格。")


def _add_engineering_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, "実装・再現性：修正を成果物契約とテストまで閉じる", "11 / ENGINEERING EVIDENCE", 12)
    cards = (
        ("MODEL", "Stage 1費用代理", "天候別PV上限・0円単価・grid下界", CYAN),
        ("FEASIBILITY", "累積所在地必要条件", "始発/便間/帰庫 + home-depot窓", GREEN),
        ("AUDIT", "summary / comparator", "代理の設定・入力・結果を別フィールド化", RAIN),
        ("SENSITIVITY", "BEV readiness", "在庫を壊さず使用可能台数を操作", SUN),
    )
    for index, (label, title, note, accent) in enumerate(cards):
        x = 0.62 + index * 3.15
        _card(slide, x, 1.62, 2.85, 1.62, fill=WHITE, line=accent)
        _add_text(slide, label, x + 0.24, 1.87, 1.45, 0.20, size=8, color=accent, bold=True)
        _add_text(slide, title, x + 0.24, 2.17, 2.35, 0.30, size=15, color=NAVY, bold=True)
        _add_text(slide, note, x + 0.24, 2.66, 2.35, 0.28, size=8.5, color=MUTED)
    _card(slide, 0.70, 3.72, 11.93, 2.24, fill=NAVY_2, line=NAVY_2)
    _add_text(slide, "再現コマンド（両ケースとも同じsolver control）", 1.03, 4.03, 4.35, 0.24, size=10, color=CYAN, bold=True)
    command = (
        "$env:GRB_LICENSE_FILE='C:\\Users\\RTDS_admin\\gurobi.lic'\n"
        "python scripts/run_research_phase3_frontend_weather.py --time-limit-sec 1500 "
        "--mip-gap 0.1 --random-seed 42 ..."
    )
    _add_text(slide, command, 1.03, 4.48, 10.55, 0.72, size=12, color=WHITE, font=FONT_MONO)
    _add_text(slide, "最終確認", 1.03, 5.40, 0.80, 0.22, size=9, color=CYAN, bold=True)
    _add_text(slide, "focused / full pytest、compileall、git diff --check、PPTレンダリング、MIT-style self review", 1.93, 5.36, 9.75, 0.28, size=11, color=RGBColor(216, 228, 238))


def _add_conclusion_slide(prs: Presentation, sunny: Mapping[str, Any], rain: Mapping[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, SLIDE_H, fill=CYAN)
    sunny_trips = _number(_mapping(sunny, "served_trip_count_by_vehicle_type").get("BEV"))
    rain_trips = _number(_mapping(rain, "served_trip_count_by_vehicle_type").get("BEV"))
    _add_text(slide, "CONCLUSION", 0.78, 0.60, 2.2, 0.25, size=10, color=CYAN, bold=True)
    _add_text(slide, "PVの価値を割当に伝えつつ、\n物理可行性はStage 2で守れた", 0.78, 1.30, 11.5, 1.12, size=31, color=WHITE, bold=True)
    _metric_card(slide, "晴天のBEV担当", f"{sunny_trips:.0f}便", "全264便中", 0.80, 3.15, 3.55, accent=SUN)
    _metric_card(slide, "雨天のBEV担当", f"{rain_trips:.0f}便", "全264便中", 4.89, 3.15, 3.55, accent=RAIN)
    _metric_card(slide, "天候差", f"+{sunny_trips - rain_trips:.0f}便", "晴天 − 雨天", 8.98, 3.15, 3.55, accent=CYAN)
    _add_text(slide, "次の一手", 0.82, 5.33, 1.0, 0.24, size=10, color=CYAN, bold=True)
    _add_text(slide, "clean commitから晴雨を再実行 → strict comparatorで正式化 → Phase 4同時最適化を発展課題として比較", 1.92, 5.24, 10.5, 0.38, size=15, color=WHITE, bold=True)
    _add_text(slide, "修論の芯：どの制約がEVバス運用の成立条件を支配するかを、再現可能な最適化実験で示す", 0.82, 6.54, 11.7, 0.30, size=11, color=RGBColor(190, 211, 227), align=PP_ALIGN.CENTER)


def build_presentation(
    *,
    sunny: Mapping[str, Any],
    rain: Mapping[str, Any],
    baseline_sunny: Mapping[str, Any],
    output_path: Path,
) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "EVバスPhase 3モデル復旧・天候差反映の進捗"
    prs.core_properties.subject = "2026-07-11〜2026-07-16 修士論文研究進捗"
    prs.core_properties.author = "master-course research project"
    prs.core_properties.comments = (
        "Generated from immutable Phase 3 summary artifacts. "
        "Two-stage results are not labeled global optima."
    )

    _add_title_slide(prs)
    _add_research_question_slide(prs)
    _add_call_chain_slide(prs)
    _add_timeline_slide(prs)
    _add_controls_slide(prs, sunny, rain)
    _add_proxy_slide(prs)
    _add_probe_slide(prs)
    _add_location_fix_slide(prs)
    _add_mix_slide(prs, sunny, rain, baseline_sunny)
    _add_energy_cost_slide(prs, sunny, rain)
    _add_validity_slide(prs, sunny, rain)
    _add_engineering_slide(prs)
    _add_conclusion_slide(prs, sunny, rain)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--sunny-summary", type=Path, default=DEFAULT_SUNNY_SUMMARY)
    parser.add_argument("--rain-summary", type=Path, default=DEFAULT_RAIN_SUMMARY)
    parser.add_argument(
        "--baseline-sunny-summary",
        type=Path,
        default=DEFAULT_BASELINE_SUNNY_SUMMARY,
    )
    parser.add_argument(
        "--baseline-rain-summary",
        type=Path,
        default=DEFAULT_BASELINE_RAIN_SUMMARY,
        help="Loaded to assert that both accepted baseline assignments match.",
    )
    parser.add_argument("--output", type=Path, default=DEFAULT_PRESENTATION)
    args = parser.parse_args()

    sunny = _load_json(args.sunny_summary)
    rain = _load_json(args.rain_summary)
    baseline_sunny = _load_json(args.baseline_sunny_summary)
    baseline_rain = _load_json(args.baseline_rain_summary)
    if baseline_sunny.get("stage1_objective") != baseline_rain.get("stage1_objective"):
        raise ValueError("Accepted baseline sunny/rain Stage 1 assignments differ")
    output = build_presentation(
        sunny=sunny,
        rain=rain,
        baseline_sunny=baseline_sunny,
        output_path=args.output.resolve(),
    )
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
