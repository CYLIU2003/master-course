"""Build the advisor-facing Phase 3 weather and energy-balance progress deck.

The deck is generated from the immutable sunny/rain audit artifacts.  Its
visual language follows the supplied September presentation: white canvas,
blue headings and rules, Meiryo typography, prominent comparison figures, and
short bottom-line conclusions.  Every quantitative chart contains both
weather cases, and every slide has a Japanese speaker script in its notes.
"""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any, Mapping, Sequence

import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_AUDIT = Path(
    r"C:\master-course\output\phase3_weather_energy_audit_20260716\weather_energy_balance_audit.json"
)
DEFAULT_PRESENTATION = (
    REPO_ROOT
    / "docs"
    / "presentations"
    / "phase3_weather_energy_balance_progress_20260716.pptx"
)
DEFAULT_ASSET_DIR = Path(
    r"C:\master-course\output\phase3_weather_energy_audit_20260716\presentation_assets"
)

SLIDE_W = 13.333
SLIDE_H = 7.5
TOTAL_SLIDES = 18
FONT = "Meiryo"

BLUE = RGBColor(0, 92, 170)
BLUE_DARK = RGBColor(0, 54, 103)
BLUE_LIGHT = RGBColor(225, 239, 250)
SUN = RGBColor(233, 154, 35)
SUN_LIGHT = RGBColor(255, 244, 219)
RAIN = RGBColor(63, 109, 180)
RAIN_LIGHT = RGBColor(229, 237, 249)
GREEN = RGBColor(35, 142, 94)
GREEN_LIGHT = RGBColor(227, 246, 237)
RED = RGBColor(194, 55, 55)
RED_LIGHT = RGBColor(252, 232, 232)
GRAY = RGBColor(92, 102, 112)
GRAY_LIGHT = RGBColor(242, 244, 246)
BLACK = RGBColor(25, 31, 36)
WHITE = RGBColor(255, 255, 255)

MPL_SUN = "#E99A23"
MPL_RAIN = "#3F6DB4"
MPL_BLUE = "#005CAA"
MPL_GREEN = "#238E5E"
MPL_TEAL = "#2A9D8F"
MPL_RED = "#C23737"
MPL_GRAY = "#5C6670"
MPL_LIGHT_BLUE = "#9EC9EA"


def _load_json(path: Path) -> dict[str, Any]:
    if not path.is_file():
        raise FileNotFoundError(f"Required audit artifact is missing: {path}")
    return json.loads(path.read_text(encoding="utf-8"))


def _number(value: Any, default: float = 0.0) -> float:
    try:
        parsed = float(value)
    except (TypeError, ValueError):
        return default
    return parsed if math.isfinite(parsed) else default


def _configure_matplotlib() -> None:
    plt.rcParams.update(
        {
            "font.family": FONT,
            "axes.unicode_minus": False,
            "axes.edgecolor": "#AAB3BB",
            "axes.labelcolor": "#25313A",
            "xtick.color": "#4D5963",
            "ytick.color": "#4D5963",
            "grid.color": "#D8DEE4",
            "grid.linewidth": 0.7,
            "axes.titleweight": "bold",
            "axes.titlesize": 12,
            "axes.labelsize": 10,
            "legend.fontsize": 9,
        }
    )


def _save_figure(fig: Any, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _clock_ticks(rows: Sequence[Mapping[str, Any]]) -> tuple[list[int], list[str]]:
    positions = list(range(0, len(rows), 4))
    labels = [str(rows[index]["slot_label"]).split("–", 1)[0] for index in positions]
    return positions, labels


def _generate_vehicle_chart(audit: Mapping[str, Any], path: Path) -> None:
    sunny = audit["cases"]["sunny"]
    rain = audit["cases"]["rain"]
    fig, axes = plt.subplots(1, 2, figsize=(10.8, 3.55), gridspec_kw={"wspace": 0.30})
    x = [0, 1]
    width = 0.34
    sunny_used = [sunny["operation"]["used_vehicle_count"]["BEV"], sunny["operation"]["used_vehicle_count"]["ICE"]]
    rain_used = [rain["operation"]["used_vehicle_count"]["BEV"], rain["operation"]["used_vehicle_count"]["ICE"]]
    axes[0].bar([value - width / 2 for value in x], sunny_used, width, label="晴天", color=MPL_SUN)
    axes[0].bar([value + width / 2 for value in x], rain_used, width, label="雨天", color=MPL_RAIN)
    for index, fleet_key in enumerate(("BEV", "ICE")):
        fleet = int(sunny["fleet_input"][fleet_key])
        axes[0].plot([index - 0.30, index + 0.30], [fleet, fleet], color=MPL_RED, linestyle="--", linewidth=1.5)
        axes[0].text(index, fleet + 0.7, f"在庫 {fleet}台", ha="center", color=MPL_RED, fontsize=9)
    for bars in axes[0].containers:
        axes[0].bar_label(bars, fmt="%.0f台", padding=2, fontsize=9)
    axes[0].set_xticks(x, ["EV", "エンジン"])
    axes[0].set_ylabel("使用車両数 [台]")
    axes[0].set_title("使用車両数：在庫全数は使っていない")
    axes[0].set_ylim(0, 39)
    axes[0].grid(axis="y")
    axes[0].legend(frameon=False, loc="upper left")

    cases = ["晴天", "雨天"]
    bev_trips = [sunny["operation"]["assigned_trip_count"]["BEV"], rain["operation"]["assigned_trip_count"]["BEV"]]
    ice_trips = [sunny["operation"]["assigned_trip_count"]["ICE"], rain["operation"]["assigned_trip_count"]["ICE"]]
    axes[1].bar(cases, bev_trips, label="EV便", color=MPL_BLUE)
    axes[1].bar(cases, ice_trips, bottom=bev_trips, label="エンジン便", color=MPL_GRAY)
    for index, (bev, ice) in enumerate(zip(bev_trips, ice_trips)):
        axes[1].text(index, bev / 2, f"EV {bev}便", ha="center", va="center", color="white", fontweight="bold")
        axes[1].text(index, bev + ice / 2, f"ICE {ice}便", ha="center", va="center", color="white", fontweight="bold")
    axes[1].set_ylabel("担当便数 [便/日]")
    axes[1].set_title("264便の割当構成")
    axes[1].set_ylim(0, 285)
    axes[1].grid(axis="y")
    axes[1].legend(frameon=False, loc="upper right")
    _save_figure(fig, path)


def _generate_active_operation_chart(audit: Mapping[str, Any], path: Path) -> None:
    fig, axes = plt.subplots(2, 1, figsize=(10.8, 4.2), sharex=True, sharey=True)
    for axis, case_key, color in zip(axes, ("sunny", "rain"), (MPL_SUN, MPL_RAIN)):
        case = audit["cases"][case_key]
        rows = case["slot_rows"]
        x = list(range(len(rows)))
        bev = [row["active_bev_count"] for row in rows]
        ice = [row["active_ice_count"] for row in rows]
        axis.step(x, bev, where="mid", color=MPL_BLUE, linewidth=2.1, label="稼働中EV")
        axis.step(x, ice, where="mid", color=MPL_GRAY, linewidth=2.1, label="稼働中エンジン")
        axis.fill_between(x, 0, bev, step="mid", color=MPL_LIGHT_BLUE, alpha=0.28)
        axis.set_title(f"{case['case_label']}（{case['service_date']}）", loc="left", color=color)
        axis.set_ylabel("稼働台数 [台]")
        axis.grid(axis="y")
        axis.yaxis.set_major_locator(MaxNLocator(integer=True))
    ticks, labels = _clock_ticks(audit["cases"]["sunny"]["slot_rows"])
    axes[-1].set_xticks(ticks, labels)
    axes[-1].set_xlabel("時刻（各60分枠の開始）")
    axes[0].legend(frameon=False, loc="upper right", ncol=2)
    fig.suptitle("時刻別の営業運行車両数（両天候を同一軸で比較）", fontsize=13, fontweight="bold")
    _save_figure(fig, path)


def _generate_bess_soc_chart(audit: Mapping[str, Any], path: Path) -> None:
    fig, axis = plt.subplots(figsize=(10.8, 3.7))
    for case_key, color, marker in (("sunny", MPL_SUN, "o"), ("rain", MPL_RAIN, "s")):
        case = audit["cases"][case_key]
        rows = case["slot_rows"]
        values = [case["bess"]["initial_soc_kwh"]] + [row["bess_soc_end_kwh"] for row in rows]
        axis.plot(range(len(values)), values, color=color, linewidth=2.5, marker=marker, markevery=4, label=case["case_label"])
    sunny = audit["cases"]["sunny"]
    axis.axhspan(sunny["bess"]["soc_min_kwh"], sunny["bess"]["soc_max_kwh"], color="#DDEFE6", alpha=0.45, label="許容SOC範囲")
    axis.axhline(sunny["bess"]["terminal_target_kwh"], color=MPL_GREEN, linestyle="--", linewidth=1.7, label="初期・終端目標 300 kWh")
    ticks = list(range(0, 25, 4))
    start_hour = 5
    labels = [f"{(start_hour + value) % 24:02d}:00" for value in ticks]
    axis.set_xticks(ticks, labels)
    axis.set_ylabel("BESS SOC [kWh]")
    axis.set_xlabel("時刻")
    axis.set_ylim(95, 505)
    axis.set_title("定置蓄電池SOC：晴天・雨天とも日末に初期値へ復帰")
    axis.grid(axis="y")
    axis.legend(frameon=False, loc="lower center", ncol=4, bbox_to_anchor=(0.5, -0.33))
    _save_figure(fig, path)


def _generate_bess_flow_chart(audit: Mapping[str, Any], path: Path) -> None:
    fig, axes = plt.subplots(2, 1, figsize=(10.8, 4.25), sharex=True, sharey=True)
    for axis, case_key, title_color in zip(axes, ("sunny", "rain"), (MPL_SUN, MPL_RAIN)):
        case = audit["cases"][case_key]
        rows = case["slot_rows"]
        x = list(range(len(rows)))
        pv_charge = [row["pv_to_bess_kwh"] for row in rows]
        grid_charge = [row["grid_to_bess_kwh"] for row in rows]
        discharge = [-row["bess_to_bus_kwh"] for row in rows]
        axis.bar(x, pv_charge, color=MPL_GREEN, label="PV→BESS（入力）")
        axis.bar(x, grid_charge, bottom=pv_charge, color=MPL_BLUE, label="系統→BESS（入力）")
        axis.bar(x, discharge, color=MPL_RED, label="BESS→バス（供給、負表示）")
        axis.axhline(0, color="#68727B", linewidth=0.8)
        axis.set_title(f"{case['case_label']}（{case['service_date']}）", loc="left", color=title_color)
        axis.set_ylabel("1時間量 [kWh]")
        axis.grid(axis="y")
    ticks, labels = _clock_ticks(audit["cases"]["sunny"]["slot_rows"])
    axes[-1].set_xticks(ticks, labels)
    axes[-1].set_xlabel("時刻（正：充電入力、負：バスへの放電供給）")
    axes[0].legend(frameon=False, loc="upper right", ncol=3)
    fig.suptitle("BESS充放電：系統充電は禁止、PVだけを蓄電", fontsize=13, fontweight="bold")
    _save_figure(fig, path)


def _generate_pv_chart(audit: Mapping[str, Any], path: Path) -> None:
    fig, axes = plt.subplots(2, 1, figsize=(10.8, 4.25), sharex=True, sharey=True)
    for axis, case_key, title_color in zip(axes, ("sunny", "rain"), (MPL_SUN, MPL_RAIN)):
        case = audit["cases"][case_key]
        rows = case["slot_rows"]
        x = list(range(len(rows)))
        direct = [row["pv_to_bus_kwh"] for row in rows]
        to_bess = [row["pv_to_bess_kwh"] for row in rows]
        curtail = [row["pv_curtail_kwh"] for row in rows]
        generation = [row["pv_generated_kwh"] for row in rows]
        axis.bar(x, direct, color=MPL_BLUE, label="PV→バス")
        axis.bar(x, to_bess, bottom=direct, color=MPL_GREEN, label="PV→BESS")
        bottom = [left + middle for left, middle in zip(direct, to_bess)]
        axis.bar(x, curtail, bottom=bottom, color="#D9A441", label="出力抑制")
        axis.plot(x, generation, color=MPL_RED, linewidth=1.8, marker="o", markersize=3, label="PV発電")
        axis.set_title(f"{case['case_label']}（{case['service_date']}）", loc="left", color=title_color)
        axis.set_ylabel("1時間量 [kWh]")
        axis.grid(axis="y")
    ticks, labels = _clock_ticks(audit["cases"]["sunny"]["slot_rows"])
    axes[-1].set_xticks(ticks, labels)
    axes[-1].set_xlabel("時刻")
    axes[0].legend(frameon=False, loc="upper right", ncol=4)
    fig.suptitle("PV発電の行き先：直接充電・BESS充電・出力抑制", fontsize=13, fontweight="bold")
    _save_figure(fig, path)


def _generate_bus_source_chart(audit: Mapping[str, Any], path: Path) -> None:
    fig, axes = plt.subplots(2, 1, figsize=(10.8, 4.25), sharex=True, sharey=True)
    for axis, case_key, title_color in zip(axes, ("sunny", "rain"), (MPL_SUN, MPL_RAIN)):
        case = audit["cases"][case_key]
        rows = case["slot_rows"]
        x = list(range(len(rows)))
        grid = [row["grid_to_bus_kwh"] for row in rows]
        pv = [row["pv_to_bus_kwh"] for row in rows]
        bess = [row["bess_to_bus_kwh"] for row in rows]
        total = [row["bus_charging_input_kwh"] for row in rows]
        axis.bar(x, grid, color=MPL_BLUE, label="系統→バス")
        axis.bar(x, pv, bottom=grid, color=MPL_GREEN, label="PV→バス")
        bottom = [left + middle for left, middle in zip(grid, pv)]
        axis.bar(x, bess, bottom=bottom, color=MPL_RED, label="BESS→バス")
        axis.plot(x, total, color="#20262B", linewidth=1.6, marker="o", markersize=3, label="充電入力合計")
        axis.set_title(f"{case['case_label']}（{case['service_date']}）", loc="left", color=title_color)
        axis.set_ylabel("1時間量 [kWh]")
        axis.grid(axis="y")
    ticks, labels = _clock_ticks(audit["cases"]["sunny"]["slot_rows"])
    axes[-1].set_xticks(ticks, labels)
    axes[-1].set_xlabel("時刻")
    axes[0].legend(frameon=False, loc="upper right", ncol=4)
    fig.suptitle("EVバス充電の電力源：系統・PV・BESSの実配分", fontsize=13, fontweight="bold")
    _save_figure(fig, path)


def _generate_grid_chart(audit: Mapping[str, Any], path: Path) -> None:
    fig, axis = plt.subplots(figsize=(10.8, 3.7))
    for case_key, color, marker in (("sunny", MPL_SUN, "o"), ("rain", MPL_RAIN, "s")):
        case = audit["cases"][case_key]
        values = [row["grid_import_average_kw"] for row in case["slot_rows"]]
        x = list(range(len(values)))
        axis.step(x, values, where="mid", color=color, linewidth=2.2, label=f"{case['case_label']}（日量 {case['daily_energy']['grid_import_kwh']:.1f} kWh）")
        peak_index = max(range(len(values)), key=lambda index: values[index])
        axis.scatter([peak_index], [values[peak_index]], s=52, color=color, marker=marker, zorder=3)
        axis.annotate(f"{values[peak_index]:.1f} kW", (peak_index, values[peak_index]), xytext=(5, 7), textcoords="offset points", color=color, fontweight="bold")
    ticks, labels = _clock_ticks(audit["cases"]["sunny"]["slot_rows"])
    axis.set_xticks(ticks, labels)
    axis.set_xlabel("時刻（60分平均）")
    axis.set_ylabel("系統購入電力 [kW]")
    axis.set_title("系統購入：雨天は日量が小さくてもピークが高い")
    axis.grid(axis="y")
    axis.legend(frameon=False, loc="upper left")
    _save_figure(fig, path)


def _generate_fuel_chart(audit: Mapping[str, Any], path: Path) -> None:
    sunny = audit["cases"]["sunny"]
    rain = audit["cases"]["rain"]
    fig, axes = plt.subplots(1, 2, figsize=(10.8, 3.6), gridspec_kw={"wspace": 0.30})
    labels = ["晴天", "雨天"]
    service = [sunny["fuel"]["service_distance_km"], rain["fuel"]["service_distance_km"]]
    deadhead = [sunny["fuel"]["intertrip_deadhead_distance_km"], rain["fuel"]["intertrip_deadhead_distance_km"]]
    axes[0].bar(labels, service, color=MPL_BLUE, label="営業距離")
    axes[0].bar(labels, deadhead, bottom=service, color=MPL_GRAY, label="便間回送距離")
    for index, (service_km, deadhead_km) in enumerate(zip(service, deadhead)):
        axes[0].text(index, service_km + deadhead_km + 30, f"{service_km + deadhead_km:,.1f} km", ha="center", fontweight="bold")
    axes[0].set_ylabel("ICE担当距離 [km/日]")
    axes[0].set_title("エンジンバス担当距離")
    axes[0].set_ylim(0, 1750)
    axes[0].grid(axis="y")
    axes[0].legend(frameon=False, loc="upper left")

    rows_s = sunny["slot_rows"]
    rows_r = rain["slot_rows"]
    x = list(range(len(rows_s)))
    fuel_s = [row["fuel_service_l"] + row["fuel_intertrip_deadhead_l"] for row in rows_s]
    fuel_r = [row["fuel_service_l"] + row["fuel_intertrip_deadhead_l"] for row in rows_r]
    axes[1].plot(x, fuel_s, color=MPL_SUN, marker="o", markersize=3, linewidth=2, label=f"晴天 {sunny['fuel']['total_fuel_l']:.1f} L")
    axes[1].plot(x, fuel_r, color=MPL_RAIN, marker="s", markersize=3, linewidth=2, label=f"雨天 {rain['fuel']['total_fuel_l']:.1f} L")
    ticks, tick_labels = _clock_ticks(rows_s)
    axes[1].set_xticks(ticks, tick_labels)
    axes[1].set_xlabel("便の出発時刻")
    axes[1].set_ylabel("距離ベース燃料 [L/時間枠]")
    axes[1].set_title("運行割当から再計算した燃料量")
    axes[1].grid(axis="y")
    axes[1].legend(frameon=False, loc="upper left")
    _save_figure(fig, path)


def _generate_cost_chart(audit: Mapping[str, Any], path: Path) -> None:
    sunny = audit["cases"]["sunny"]
    rain = audit["cases"]["rain"]
    labels = ["晴天", "雨天"]
    component_specs = [
        ("vehicle_usage_cost", "車両使用", "#B7BEC5"),
        ("electricity_cost", "電力", MPL_BLUE),
        ("fuel_cost", "燃料", MPL_GRAY),
        ("demand_cost", "需要料金", MPL_RED),
        ("co2_cost", "CO₂", MPL_GREEN),
    ]
    fig, axes = plt.subplots(1, 2, figsize=(10.8, 3.65), gridspec_kw={"wspace": 0.32})
    bottoms = [0.0, 0.0]
    for key, label, color in component_specs:
        values = [sunny["costs_jpy"][key] / 1000.0, rain["costs_jpy"][key] / 1000.0]
        axes[0].bar(labels, values, bottom=bottoms, label=label, color=color)
        bottoms = [bottom + value for bottom, value in zip(bottoms, values)]
    for index, total in enumerate(bottoms):
        axes[0].text(index, total + 12, f"{total:,.1f}千円", ha="center", fontweight="bold")
    axes[0].set_ylabel("会計総額 [千円/日]")
    axes[0].set_title("総額（車両使用費を含む）")
    axes[0].set_ylim(0, 790)
    axes[0].grid(axis="y")

    x = [0, 1]
    width = 0.18
    variable_specs = component_specs[1:]
    for offset_index, (key, label, color) in enumerate(variable_specs):
        values = [sunny["costs_jpy"][key] / 1000.0, rain["costs_jpy"][key] / 1000.0]
        offsets = [value + (offset_index - 1.5) * width for value in x]
        axes[1].bar(offsets, values, width, label=label, color=color)
    axes[1].set_xticks(x, labels)
    axes[1].set_ylabel("費用 [千円/日]")
    axes[1].set_title("変動要因を拡大（車両使用費を除く）")
    axes[1].grid(axis="y")
    handles, legend_labels = axes[0].get_legend_handles_labels()
    fig.legend(
        handles,
        legend_labels,
        frameon=False,
        ncol=5,
        loc="lower center",
        bbox_to_anchor=(0.5, -0.07),
    )
    _save_figure(fig, path)


def _set_background(slide: Any, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: RGBColor = BLACK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _add_rich_text(
    slide: Any,
    runs: Sequence[tuple[str, float, RGBColor, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, size, color, bold in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return shape


def _add_rect(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius: bool = False,
) -> Any:
    # Keep the visual language close to the supplied academic reference deck:
    # square-edged panels, restrained fills, and no app-style rounded cards.
    # ``radius`` remains in the public helper signature for call-site stability.
    del radius
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _add_tcu_mark(slide: Any) -> None:
    """Draw a stable vector approximation of the supplied TCU deck mark."""

    cyan = RGBColor(0, 164, 222)
    _add_rect(slide, 11.58, 0.15, 0.31, 0.18, fill=cyan, line=cyan)
    _add_rect(slide, 11.64, 0.38, 0.08, 0.08, fill=cyan, line=cyan)
    _add_text(
        slide,
        "東京都市大学",
        11.95,
        0.10,
        1.08,
        0.25,
        size=10.5,
        color=BLACK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.0,
    )
    _add_text(
        slide,
        "TOKYO CITY UNIVERSITY",
        11.95,
        0.35,
        1.08,
        0.12,
        size=4.7,
        color=BLACK,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.0,
    )


def _add_card(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    title: str,
    value: str,
    detail: str,
    accent: RGBColor,
    fill: RGBColor,
) -> None:
    _add_rect(slide, x, y, w, h, fill=fill, line=accent, radius=True)
    _add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.35, size=13, color=accent, bold=True)
    _add_text(slide, value, x + 0.18, y + 0.52, w - 0.36, 0.62, size=25, color=BLACK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _add_text(slide, detail, x + 0.18, y + 1.16, w - 0.36, h - 1.28, size=11.5, color=GRAY)


def _add_header(slide: Any, title: str, page: int) -> None:
    _set_background(slide)
    _add_text(slide, title, 0.30, 0.12, 10.7, 0.48, size=24, color=BLUE_DARK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.67), Inches(SLIDE_W), Inches(0.67))
    line.line.color.rgb = BLUE
    line.line.width = Pt(2.0)
    _add_tcu_mark(slide)
    _add_text(slide, f"EVバスPhase 3進捗　2026-07-16　{page}/{TOTAL_SLIDES}", 0.30, 7.18, 12.7, 0.22, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT)


def _add_bottom_line(slide: Any, text: str, *, fill: RGBColor = BLUE_LIGHT, color: RGBColor = BLUE_DARK) -> None:
    _add_rect(slide, 0.35, 6.50, 12.63, 0.54, fill=fill, line=fill, radius=True)
    _add_text(slide, f"要点｜{text}", 0.55, 6.60, 12.20, 0.30, size=14.5, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)


def _add_picture(slide: Any, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _add_notes(slide: Any, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _add_bullets(
    slide: Any,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: RGBColor = BLACK,
) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"● {item}"
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(9)


def _add_table(
    slide: Any,
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    column_widths: Sequence[float],
    font_size: float = 12.5,
    header_font_size: float = 13.0,
) -> Any:
    """Add a restrained research-style table with alternating row fills."""

    if not rows or any(len(row) != len(rows[0]) for row in rows):
        raise ValueError("table rows must be non-empty and rectangular")
    if len(column_widths) != len(rows[0]):
        raise ValueError("column width count must match table column count")
    if abs(sum(column_widths) - w) > 0.02:
        raise ValueError("column widths must add up to table width")

    table_shape = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = table_shape.table
    for index, width in enumerate(column_widths):
        table.columns[index].width = Inches(width)
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                BLUE_DARK if row_index == 0 else (GRAY_LIGHT if row_index % 2 == 0 else WHITE)
            )
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if column_index else PP_ALIGN.LEFT
                paragraph.font.name = FONT
                paragraph.font.size = Pt(header_font_size if row_index == 0 else font_size)
                paragraph.font.bold = row_index == 0 or column_index == 0
                paragraph.font.color.rgb = WHITE if row_index == 0 else BLACK
    return table_shape


def _add_flow_column(slide: Any, case: Mapping[str, Any], x: float, accent: RGBColor, fill: RGBColor) -> None:
    daily = case["daily_energy"]
    bess = case["bess"]
    _add_rect(slide, x, 1.05, 5.95, 4.98, fill=WHITE, line=accent, radius=True)
    _add_text(slide, f"{case['case_label']}  {case['service_date']}", x + 0.18, 1.15, 5.58, 0.38, size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, x + 0.25, 1.70, 1.35, 0.68, fill=fill, line=accent, radius=True)
    _add_text(slide, f"PV\n{daily['pv_generated_kwh']:.1f}", x + 0.30, 1.78, 1.25, 0.50, size=13, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _add_text(slide, "→", x + 1.65, 1.88, 0.34, 0.30, size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    _add_text(
        slide,
        f"バス {daily['pv_to_bus_kwh']:.1f}\nBESS {daily['pv_to_bess_kwh']:.1f}\n抑制 {daily['pv_curtailed_kwh']:.1f} kWh",
        x + 2.02,
        1.66,
        1.72,
        0.82,
        size=12.5,
        color=BLACK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _add_rect(slide, x + 3.96, 1.70, 1.65, 0.68, fill=BLUE_LIGHT, line=BLUE, radius=True)
    _add_text(slide, f"系統購入\n{daily['grid_import_kwh']:.1f} kWh", x + 4.04, 1.78, 1.50, 0.50, size=12.5, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _add_rect(slide, x + 0.25, 2.72, 2.30, 1.10, fill=GREEN_LIGHT, line=GREEN, radius=True)
    _add_text(slide, "BESS", x + 0.38, 2.91, 0.72, 0.28, size=14, color=GREEN, bold=True)
    _add_text(
        slide,
        f"入力 {daily['bess_charge_input_kwh']:.1f}\nバス供給 {daily['bess_discharge_delivered_kwh']:.1f} kWh",
        x + 1.18,
        2.79,
        1.23,
        0.78,
        size=11.5,
        color=BLACK,
        bold=True,
    )
    _add_text(slide, "→", x + 2.62, 3.08, 0.34, 0.30, size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, x + 3.02, 2.72, 2.59, 1.10, fill=GRAY_LIGHT, line=GRAY, radius=True)
    _add_text(
        slide,
        f"充電入力合計\n{daily['bus_charging_input_kwh']:.1f} kWh\n= 系統 + PV + BESS",
        x + 3.12,
        2.80,
        2.38,
        0.82,
        size=12,
        color=BLACK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    total_loss = daily["bess_charge_loss_kwh"] + daily["bess_discharge_loss_kwh"]
    residual = max(case["balances"]["max_absolute_residuals"].values())
    _add_text(
        slide,
        f"BESS: {bess['initial_soc_kwh']:.0f} → {bess['terminal_soc_kwh']:.0f} kWh　"
        f"往復損失 {total_loss:.1f} kWh",
        x + 0.28,
        4.13,
        5.34,
        0.38,
        size=13.5,
        color=BLUE_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        f"PV・充電源・BESS式の最大残差\n{residual:.2e} kWh  ≤  1.0e-6 kWh",
        x + 0.53,
        4.70,
        4.84,
        0.72,
        size=14,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build_deck(audit: Mapping[str, Any], output_path: Path, asset_dir: Path) -> None:
    _configure_matplotlib()
    asset_dir.mkdir(parents=True, exist_ok=True)
    charts = {
        "vehicles": asset_dir / "vehicle_use_comparison.png",
        "active": asset_dir / "active_operation_comparison.png",
        "bess_soc": asset_dir / "bess_soc_comparison.png",
        "bess_flow": asset_dir / "bess_flow_comparison.png",
        "pv": asset_dir / "pv_destinations_comparison.png",
        "sources": asset_dir / "bus_charging_sources_comparison.png",
        "grid": asset_dir / "grid_import_comparison.png",
        "fuel": asset_dir / "fuel_operation_comparison.png",
        "cost": asset_dir / "cost_comparison.png",
    }
    _generate_vehicle_chart(audit, charts["vehicles"])
    _generate_active_operation_chart(audit, charts["active"])
    _generate_bess_soc_chart(audit, charts["bess_soc"])
    _generate_bess_flow_chart(audit, charts["bess_flow"])
    _generate_pv_chart(audit, charts["pv"])
    _generate_bus_source_chart(audit, charts["sources"])
    _generate_grid_chart(audit, charts["grid"])
    _generate_fuel_chart(audit, charts["fuel"])
    _generate_cost_chart(audit, charts["cost"])

    sunny = audit["cases"]["sunny"]
    rain = audit["cases"]["rain"]
    params = sunny["scenario_parameters"]
    rain_params = rain["scenario_parameters"]
    depot_assets = params["depot_energy_assets"]["tsurumaki"]
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    _set_background(slide)
    _add_tcu_mark(slide)
    _add_rect(slide, 0.40, 2.00, 12.20, 1.62, fill=BLUE_DARK, line=BLUE_DARK)
    _add_text(slide, "EVバス運用最適化 Phase 3", 0.72, 2.22, 11.56, 0.52, size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _add_text(slide, "モデル修正・実験条件・晴天／雨天比較", 0.72, 2.82, 11.56, 0.44, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "2026年7月16日　進捗報告", 0.75, 4.18, 11.50, 0.40, size=19, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "刘 承洋　東京都市大学", 0.75, 4.65, 11.50, 0.35, size=16, color=BLACK, align=PP_ALIGN.CENTER)
    _add_text(slide, "比較日：晴天 2025-08-05 ／ 雨天 2025-08-10", 0.75, 5.45, 11.50, 0.34, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    _add_notes(
        slide,
        "目標25秒\n本日は、Phase 3の最適化システムに加えた修正、実際に計算へ投入した時間・費用・設備条件、そのうえで得られた晴天と雨天の比較結果を報告します。先生から確認を求められた蓄電池の日末SOC、EV使用台数、電力需給、燃料整合も結果スライドで順に示します。数値は7月16日時点の計算成果物と保存されたsolver入力から再監査しました。",
    )

    # 2. Direct answers
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "結論：先生の確認事項への回答", 2)
    direct_answer_rows = [
        ["確認項目", "晴天 2025-08-05", "雨天 2025-08-10", "判定・解釈"],
        ["BESS日末SOC", "300 → 300 kWh", "300 → 300 kWh", "終端300 kWhのハード制約を満足"],
        ["使用EV台数", "16 / 35台", "15 / 35台", "晴天でも全35台は使用しない"],
        ["EV担当便", "141 / 264便", "119 / 264便", "晴天は雨天より22便多い"],
        [
            "電力需給の最大残差",
            f"{max(sunny['balances']['max_absolute_residuals'].values()):.2e} kWh",
            f"{max(rain['balances']['max_absolute_residuals'].values()):.2e} kWh",
            "許容値1.0e-6 kWh以内",
        ],
        [
            "燃料使用量 / 燃料費",
            f"{sunny['fuel']['total_fuel_l']:.3f} L / {sunny['fuel']['reported_fuel_cost_final_jpy']:,.0f}円",
            f"{rain['fuel']['total_fuel_l']:.3f} L / {rain['fuel']['reported_fuel_cost_final_jpy']:,.0f}円",
            "営業＋便間回送距離による仮勘定",
        ],
    ]
    _add_table(
        slide,
        direct_answer_rows,
        0.45,
        1.03,
        12.43,
        4.98,
        column_widths=[2.55, 2.65, 2.65, 4.58],
        font_size=13.0,
        header_font_size=13.5,
    )
    _add_bottom_line(slide, "主要な整合条件は満足した。ICE在庫数と燃料台帳は正式KPI化前の確認事項として残る。")
    _add_notes(
        slide,
        "目標45秒\n主要な確認結果を表にまとめます。蓄電池は両日とも300キロワット時で始まり、300キロワット時で終了しました。晴天のEV使用は16台、雨天は15台で、晴天でも在庫35台すべては使用していません。EV担当便は晴天が141便、雨天が119便です。PV配分、バス充電源、BESS遷移の最大残差はいずれも許容値以下でした。燃料費は運行距離と一致しますが、給油イベントを含む実績台帳ではなく距離に基づく仮勘定です。",
    )

    # 3. Conditions and limitations
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "比較条件と解釈上の注意", 3)
    data = [
        ["項目", "晴天", "雨天"],
        ["日付 / scenario", f"{sunny['service_date']}\n{sunny['scenario_id'][:8]}…", f"{rain['service_date']}\n{rain['scenario_id'][:8]}…"],
        ["対象便・在庫", "264便 / EV35・ICE25", "264便 / EV35・ICE25"],
        ["PV発電", f"{sunny['daily_energy']['pv_generated_kwh']:.1f} kWh", f"{rain['daily_energy']['pv_generated_kwh']:.1f} kWh"],
        ["計算状態", f"Stage 1: time limit\ngap {sunny['solver']['stage1_mip_gap_ratio']*100:.1f}%", f"Stage 1: time limit\ngap {rain['solver']['stage1_mip_gap_ratio']*100:.1f}%"],
        ["Stage 2 / hard validation", "optimal / 全違反0", "optimal / 全違反0"],
    ]
    table_shape = slide.shapes.add_table(len(data), 3, Inches(0.55), Inches(1.05), Inches(8.15), Inches(4.95))
    table = table_shape.table
    widths = [2.25, 2.95, 2.95]
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for row_index, row in enumerate(data):
        for col_index, text in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_DARK if row_index == 0 else (GRAY_LIGHT if row_index % 2 == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_index else PP_ALIGN.LEFT
                paragraph.font.name = FONT
                paragraph.font.size = Pt(13 if row_index else 14)
                paragraph.font.bold = row_index == 0 or col_index == 0
                paragraph.font.color.rgb = WHITE if row_index == 0 else BLACK
    _add_rect(slide, 9.05, 1.05, 3.80, 2.12, fill=RED_LIGHT, line=RED, radius=True)
    _add_text(slide, "入力台数の不一致", 9.28, 1.27, 3.35, 0.35, size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "依頼文：EV35・ICE26\n現モデル：EV35・ICE25", 9.28, 1.76, 3.35, 0.75, size=18, color=BLACK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _add_text(slide, "26台条件の主張には、scenario修正後の再計算が必要。", 9.28, 2.57, 3.35, 0.38, size=11.5, color=RED, align=PP_ALIGN.CENTER)
    _add_rect(slide, 9.05, 3.47, 3.80, 2.53, fill=SUN_LIGHT, line=SUN, radius=True)
    _add_text(slide, "この資料の位置づけ", 9.28, 3.70, 3.35, 0.35, size=17, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(slide, ["二段階モデルの可行解", "総費用の大域最適解ではない", "未コミット変更を含む暫定結果", "trip/vehicle hashは再照合済み"], 9.28, 4.17, 3.25, 1.55, size=12.5)
    _add_bottom_line(slide, "以降は“暫定可行解の内部整合”として解釈する。", fill=SUN_LIGHT, color=BLUE_DARK)
    _add_notes(
        slide,
        "目標55秒\n比較条件です。両日とも264便、EV35台、ICE25台、同じ車両・便・単価・充電器・BESS条件で、主な差はPV時系列です。ここで大切なのは、先生のメモにあるICE26台と、実際にモデルへ入った25台が一致していないことです。この資料では実入力の25台として報告し、26台条件は再計算課題とします。またStage 1は制限時間終了で約13パーセントのギャップが残るため、大域最適とは言わず、Stage 2が最適かつ全ハード制約を満たす可行解として扱います。",
    )

    # 4. Optimization-system changes
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "最適化システムに加えた修正", 4)
    change_rows = [
        ["対象", "修正前に生じた問題", "今回の修正"],
        [
            "Stage 1費用評価",
            "EV充電費とPV効果を見ず、天候が変わっても同じ割当を選びやすい",
            "EV外部充電量の下界を計算し、PV・初期BESS余剰・系統の順に費用評価へ反映",
        ],
        [
            "充電場所とSOC",
            "車庫外の待機時間にも充電できる前提を暗黙に置き、Stage 2で不可行になる",
            "車庫内充電可能枠だけを使う累積SOC必要条件をStage 1へ875本追加",
        ],
        [
            "車両運用の連続性",
            "同一車両に離れた運行断片を割り当て、Stage 2で位置・SOCを接続できない",
            "1車両1日1連続仕業に限定し、Stage 1とStage 2の運行表現を一致",
        ],
        [
            "帰庫回送のSOC",
            "帰庫完了より1枠遅れて回送消費を差し引く時刻ずれ",
            "帰庫完了直後のSOC遷移で回送消費を差し引くよう修正",
        ],
        [
            "研究用実行経路",
            "fallbackや後処理がMILP出力と報告値の対応を不明瞭にする可能性",
            "fallback・postsolve repairを停止し、全便・重複・SOC・需給の検証を合格条件化",
        ],
        [
            "結果の保存",
            "PV・BESS・系統の出所と二段階モデルの意味が成果物だけでは不明確",
            "電力源別フロー、solver状態、入力hash、目的関数の意味を監査JSONへ保存",
        ],
    ]
    _add_table(
        slide,
        change_rows,
        0.52,
        0.94,
        12.29,
        5.35,
        column_widths=[2.12, 4.36, 5.81],
        font_size=11.5,
        header_font_size=13.0,
    )
    _add_bottom_line(
        slide,
        "便接続条件 arrival + turnaround + deadhead ≤ next departure は変更していない。",
        fill=GREEN_LIGHT,
        color=GREEN,
    )
    _add_notes(
        slide,
        "目標70秒\n今回の主な修正は六点です。第一に、Stage 1にもEVの外部充電費の下界を入れ、PVが多い日の効果が運行割当に反映されるようにしました。第二に、車庫外で充電したことになる不可行解を防ぐため、車庫内で利用できる充電枠だけを使う累積SOC必要条件を875本追加しました。第三に、同一車両の運行を一日一つの連続仕業に限定しました。第四に帰庫回送のSOC差引時刻を修正しました。第五に研究用実行ではfallbackと後処理修復を停止し、検証合格を必須にしました。第六に、電力源別フローとsolver情報を成果物へ残しました。便同士を接続できるかという基礎条件は変更していません。",
    )

    # 5. Two-stage model semantics
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "二段階モデルの役割と目的関数", 5)
    _add_rect(slide, 0.55, 1.02, 5.35, 2.10, fill=BLUE_LIGHT, line=BLUE)
    _add_text(slide, "Stage 1：車両運用・動力選択", 0.78, 1.20, 4.90, 0.36, size=18, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(
        slide,
        "決定：各便をどのEV／ICEへ割り当てるか\n評価：ICE燃料・CO₂、使用車両費、EV外部充電費の下界\n制約：全便、時刻重複、便接続、連続仕業、EV累積SOC必要条件",
        0.82,
        1.70,
        4.82,
        1.13,
        size=13.0,
        color=BLACK,
    )
    _add_text(slide, "→", 6.12, 1.71, 0.88, 0.48, size=34, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "割当を固定", 6.03, 2.23, 1.04, 0.28, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    _add_rect(slide, 7.28, 1.02, 5.50, 2.10, fill=GREEN_LIGHT, line=GREEN)
    _add_text(slide, "Stage 2：充電・電力需給", 7.52, 1.20, 5.02, 0.36, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_text(
        slide,
        "決定：充電時刻、系統・PV・BESS配分、系統ピーク\n評価：電力費、需要料金、CO₂、BESS費用\n制約：車両SOC、充電器、受電上限、PV収支、BESS SOC・終端条件",
        7.55,
        1.70,
        4.96,
        1.13,
        size=13.0,
        color=BLACK,
    )
    _add_rect(slide, 0.55, 3.40, 12.23, 2.58, fill=WHITE, line=GRAY)
    _add_text(slide, "Stage 1で用いるEV外部充電量の下界", 0.82, 3.62, 5.20, 0.34, size=16, color=BLUE_DARK, bold=True)
    _add_text(
        slide,
        "E_ext,v = max{0, E_trip + E_start + E_between + E_return\n　　　　　　　 + E_terminal − E_initial} / η_charge",
        0.95,
        4.07,
        5.70,
        0.86,
        size=16.5,
        color=BLACK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _add_text(
        slide,
        "下界の費用配分：PV（0円）→ 初期BESS余剰 → 最安の系統電力\n含める：天候別PV量、走行エネルギー、初期・終端SOC\n含めない：充電時刻、充電器競合、需要料金の厳密値\n　　　　　これらはStage 2で厳密に決定する",
        6.85,
        3.72,
        5.55,
        1.83,
        size=13.0,
        color=BLACK,
    )
    _add_bottom_line(
        slide,
        "本結果は二段階モデルの可行解であり、運行と充電を同時に解く総費用の大域最適解ではない。",
        fill=SUN_LIGHT,
        color=BLUE_DARK,
    )
    _add_notes(
        slide,
        "目標75秒\nモデルは二段階です。Stage 1で便と車両の割当を決め、その割当を固定してStage 2で充電と電力需給を決めます。修正前のStage 1はEV充電費とPV効果を評価していなかったため、天候差が割当に出にくい構造でした。そこで、走行、回送、終端SOCから必要な外部充電量の下界を計算し、PV、初期BESS余剰、系統の順に費用へ割り当てました。これは充電時刻や充電器競合を無視する下界です。厳密な充電実現性はStage 2で確認します。そのため、本結果を同時最適化の大域最適解とは呼びません。",
    )

    # 6. Scenario and solver parameters
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "シナリオ・計算・設備条件", 6)
    fast_chargers = sum(
        1 for charger in params["charger_configuration"] if charger["power_kw"] == 90.0
    )
    normal_chargers = sum(
        1 for charger in params["charger_configuration"] if charger["power_kw"] == 50.0
    )
    computation_rows = [
        ["計算条件", "設定値"],
        ["モデル / backend", f"{params['phase']} / {params['solver_backend']}"],
        ["総制限時間", f"{params['time_limit_sec']}秒"],
        ["段階別上限", f"Stage 1 最大{params['stage_time_limit_sec']}秒、Stage 2 最大{params['stage_time_limit_sec']}秒"],
        ["実測時間（晴天）", f"Stage 1 {sunny['solver']['stage1_runtime_seconds']:.3f}秒 / Stage 2 {sunny['solver']['stage2_runtime_seconds']:.3f}秒"],
        ["実測時間（雨天）", f"Stage 1 {rain['solver']['stage1_runtime_seconds']:.3f}秒 / Stage 2 {rain['solver']['stage2_runtime_seconds']:.3f}秒"],
        ["MIPGap / seed", f"{params['mip_gap']*100:.0f}% / {params['random_seed']}"],
        ["時間離散化", f"{params['timestep_min']}分×{params['price_slot_count']}枠、{params['horizon_start']}開始"],
        ["研究用実行", "fallbackなし、postsolve repairなし"],
    ]
    equipment_rows = [
        ["運行・設備条件", "設定値"],
        ["対象便 / 在庫", f"{params['trip_count']}便 / EV{params['fleet']['BEV']}台・ICE{params['fleet']['ICE']}台"],
        ["仕業断片", "1車両1日1連続仕業"],
        ["充電器", f"90 kW×{fast_chargers}基、50 kW×{normal_chargers}基（各1口）"],
        ["車庫受電上限", f"{params['depot_import_limit_kw_by_depot']['tsurumaki']:.0f} kW"],
        ["PV", f"{depot_assets['pv_capacity_kw']:.1f} kW：晴天{depot_assets['pv_generation_kwh']:.1f} / 雨天{rain_params['depot_energy_assets']['tsurumaki']['pv_generation_kwh']:.1f} kWh"],
        ["BESS", f"{depot_assets['bess_energy_kwh']:.0f} kWh / {depot_assets['bess_power_kw']:.0f} kW"],
        ["BESS SOC", f"初期・終端{depot_assets['bess_initial_soc_kwh']:.0f}、範囲{depot_assets['bess_soc_min_kwh']:.0f}–{depot_assets['bess_soc_max_kwh']:.0f} kWh"],
        ["BESS効率 / 経路", f"充放電各{depot_assets['bess_charge_efficiency']*100:.0f}% / PV→BESS可、系統→BESS不可"],
        ["車両終端SOC", "下限20%、目標80%±20%（実効下限60%）"],
    ]
    _add_table(slide, computation_rows, 0.48, 0.97, 6.10, 5.47, column_widths=[2.13, 3.97], font_size=11.4)
    _add_table(slide, equipment_rows, 6.75, 0.97, 6.10, 5.47, column_widths=[2.08, 4.02], font_size=11.2)
    _add_bottom_line(slide, "依頼文のICE26台に対し、保存された実入力はICE25台。26台条件の結果は別途再計算が必要。", fill=RED_LIGHT, color=RED)
    _add_notes(
        slide,
        "目標75秒\n計算条件を明記します。総制限時間は両日とも1500秒で、実装上はStage 1とStage 2へ最大750秒ずつ配分します。今回、Stage 1は両日とも約750秒を使い、晴天13.1パーセント、雨天12.9パーセントのギャップで終了しました。Stage 2は0.05秒未満で最適になりました。MIPGap設定は10パーセント、seedは42、60分24枠です。設備は90キロワット5基と50キロワット5基、受電上限1000キロワット、PV101.5キロワット、BESS600キロワット時・300キロワットです。なお保存された実入力はICE25台で、依頼文の26台とは一致していません。",
    )

    # 7. Cost and environmental parameters
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "費用・環境パラメータ", 7)
    tou_rows = [
        ["系統電力単価", "円/kWh"],
        ["05:00–15:00", "18"],
        ["16:00–17:00", "22"],
        ["18:00–23:00", "19"],
        ["00:00–04:00", "18"],
        ["需要料金", f"{params['demand_charge_monthly_yen_per_kw']:,.0f}円/kW・月"],
        ["1日換算", f"{params['demand_charge_horizon_yen_per_kw']:,.0f}円/kW"],
    ]
    other_cost_rows = [
        ["その他の単価", "設定値"],
        ["軽油", f"{params['diesel_price_yen_per_l']:,.0f}円/L"],
        ["CO₂価格", f"{params['co2_price_yen_per_kg']:.0f}円/kg-CO₂"],
        ["系統CO₂係数", "0.5 kg-CO₂/kWh"],
        ["使用車両費", f"{params['vehicle_usage_cost_jpy_per_used_bus']:,.0f}円/台・日"],
        ["PV限界単価", f"{params['pv_marginal_charge_cost_yen_per_kwh']:.0f}円/kWh"],
        ["PV抑制", f"{params['pv_curtail_penalty_yen_per_kwh']:.0f}円/kWh"],
        ["BESSサイクル費", f"{depot_assets['bess_cycle_cost_yen_per_kwh']:.0f}円/kWh"],
    ]
    _add_table(slide, tou_rows, 0.48, 1.00, 3.42, 4.76, column_widths=[2.05, 1.37], font_size=12.0)
    _add_table(slide, other_cost_rows, 4.10, 1.00, 3.75, 4.76, column_widths=[2.02, 1.73], font_size=11.8)
    _add_rect(slide, 8.08, 1.00, 4.77, 4.76, fill=WHITE, line=BLUE)
    _add_text(slide, "費用比較で用いる内訳", 8.35, 1.25, 4.23, 0.38, size=17, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(
        slide,
        "C総 = C電力 + C燃料 + C需要\n　　+ CCO₂ + C使用車両",
        8.43,
        1.84,
        4.06,
        0.88,
        size=16.5,
        color=BLACK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _add_text(
        slide,
        "有効：電力・燃料・需要料金・CO₂・使用車両費\n重み：energy=1、fuel=1、demand=1、usage=1\n　　　vehicle fixed=0、degradation=0\n\nPV単価、PV抑制、BESSサイクル費が0円なのは、今回のシナリオ設定であり、システム固有の定数ではない。",
        8.42,
        3.02,
        4.04,
        2.24,
        size=12.5,
        color=BLACK,
    )
    _add_bottom_line(slide, "需要料金は購入電力量の合計ではなく、1日の最大系統受電電力[kW]に40円/kWを乗じる。")
    _add_notes(
        slide,
        "目標70秒\n費用条件です。系統電力は5時から15時が18円、16時と17時が22円、18時から23時が19円、その後5時までは18円です。需要料金は月1200円毎キロワットを一日40円毎キロワットへ換算し、系統受電の最大値に掛けます。軽油は150円毎リットル、CO2は1円毎キログラム、系統排出係数は0.5キログラム毎キロワット時、使用車両費は一台一日2万円です。PV限界単価、PV抑制、BESSサイクル費は今回すべてゼロです。これらは今回のシナリオ条件であり、モデルが常にゼロと決めているわけではありません。",
    )

    # 8. Vehicle assignment
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "EV・エンジンバスの使用結果", 8)
    _add_picture(slide, charts["vehicles"], 0.55, 0.92, 8.55, 4.95)
    _add_rect(slide, 9.35, 1.14, 3.48, 4.42, fill=BLUE_LIGHT, line=BLUE, radius=True)
    _add_text(slide, "確認結果", 9.60, 1.38, 2.98, 0.36, size=18, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(
        slide,
        [
            "晴天：EV16台・ICE16台",
            "雨天：EV15台・ICE17台",
            "晴天EV便 141 / 264便",
            "雨天EV便 119 / 264便",
            "晴天でもEV19台が未使用",
        ],
        9.63,
        1.95,
        2.92,
        2.52,
        size=14.2,
    )
    _add_text(slide, "PV限界単価が0円でも、接続可能性・SOC・使用車両費・充電時刻の制約によりEV全車使用にはならない。", 9.63, 4.55, 2.92, 0.75, size=11.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "晴天ではEV利用が増えるが、増加は1台・22便であり全35台使用ではない。")
    _add_notes(
        slide,
        "目標50秒\n晴天ではEV16台、エンジン16台を使い、EVが141便を担当しました。雨天はEV15台、エンジン17台で、EV担当は119便です。したがって期待した方向、つまり晴天の方がEV利用が多いことは確認できます。ただし晴天でも19台のEVは未使用です。理由はPV単価だけではなく、各便の時間接続、SOC、充電可能時刻、そして一台二万円の車両使用費が同時に効くためです。",
    )

    # 9. Active operation
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "時刻別の運行状況", 9)
    _add_picture(slide, charts["active"], 0.50, 0.88, 9.60, 5.35)
    _add_rect(slide, 10.35, 1.15, 2.48, 4.70, fill=GRAY_LIGHT, line=GRAY, radius=True)
    _add_text(slide, "読み方", 10.57, 1.38, 2.04, 0.34, size=17, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(slide, ["営業運行中の実車数", "回送・充電・待機は含めない", "両日を同じ縦軸で表示", "差は便割当の時間分布にも現れる"], 10.58, 1.94, 2.00, 2.10, size=12.5)
    _add_text(slide, "日合計だけでなく、ピーク時にどの動力車が稼働するかを確認する。", 10.58, 4.55, 2.00, 0.72, size=12, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "晴天・雨天の差は、使用台数だけでなく時間帯別の動力分担として評価する。")
    _add_notes(
        slide,
        "目標35秒\nこの図は各60分枠で営業運行中の車両数です。日合計の使用台数だけでなく、いつEVとエンジン車が稼働しているかを比較しています。回送、充電、待機はこの図には含めていません。系統ピークや充電時刻を説明するときは、運行側の時間分布と組み合わせて見る必要があります。",
    )

    # 10. BESS SOC
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "蓄電池SOC：日末変化量は0", 10)
    _add_picture(slide, charts["bess_soc"], 0.55, 0.92, 9.50, 4.90)
    _add_card(slide, 10.30, 1.12, 2.55, 1.45, title="晴天", value="300 → 300", detail="範囲 120–480 kWh", accent=SUN, fill=SUN_LIGHT)
    _add_card(slide, 10.30, 2.82, 2.55, 1.45, title="雨天", value="300 → 300", detail="範囲 227.0–322.0 kWh", accent=RAIN, fill=RAIN_LIGHT)
    _add_rect(slide, 10.30, 4.54, 2.55, 1.18, fill=GREEN_LIGHT, line=GREEN, radius=True)
    _add_text(slide, "終端偏差", 10.50, 4.72, 2.15, 0.30, size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "≤ 5.7×10⁻¹⁴ kWh", 10.48, 5.08, 2.19, 0.35, size=16, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "日末SOC同値は評価結果ではなく、300 kWhの終端ハード制約により保証される。")
    _add_notes(
        slide,
        "目標45秒\n蓄電池SOCは5時の300キロワット時から始まります。晴天は朝に120まで放電し、昼のPVで480まで充電し、夜に300へ戻します。雨天はPVが少ないため変動幅が小さく、約227から322の範囲です。どちらも日末は300で、変化量は数値誤差を除きゼロです。これは偶然ではなく、終端SOC目標300キロワット時をハード制約として置いた結果です。",
    )

    # 11. BESS flows
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "蓄電池の充放電と損失", 11)
    _add_picture(slide, charts["bess_flow"], 0.48, 0.88, 9.55, 5.35)
    sunny_loss = sunny["daily_energy"]["bess_charge_loss_kwh"] + sunny["daily_energy"]["bess_discharge_loss_kwh"]
    rain_loss = rain["daily_energy"]["bess_charge_loss_kwh"] + rain["daily_energy"]["bess_discharge_loss_kwh"]
    _add_card(slide, 10.28, 1.10, 2.57, 1.55, title="晴天", value=f"入力 {sunny['daily_energy']['bess_charge_input_kwh']:.1f}", detail=f"供給342.0 / 損失 {sunny_loss:.1f} kWh", accent=SUN, fill=SUN_LIGHT)
    _add_card(slide, 10.28, 2.95, 2.57, 1.55, title="雨天", value=f"入力 {rain['daily_energy']['bess_charge_input_kwh']:.1f}", detail=f"供給90.3 / 損失 {rain_loss:.1f} kWh", accent=RAIN, fill=RAIN_LIGHT)
    _add_text(slide, "系統→BESS = 0 kWh\n（両日、設定で禁止）", 10.36, 4.90, 2.39, 0.62, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "SOC同値でも、効率95%のため入力量とバス供給量は一致せず、損失を含めて帳尻を取る。")
    _add_notes(
        slide,
        "目標45秒\n正側がBESSへの充電入力、負側がバスへの放電供給です。系統からBESSへの充電は設定で禁止されており、両日ゼロです。晴天ではPVから378.9キロワット時を入力し342.0をバスへ供給します。差はSOCに残るのではなく、充放電効率95パーセントによる往復損失36.9キロワット時です。雨天も同じ式で約9.8キロワット時の損失になります。",
    )

    # 12. PV destinations
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "PV発電出力の行き先", 12)
    _add_picture(slide, charts["pv"], 0.50, 0.88, 9.55, 5.35)
    _add_card(slide, 10.28, 1.10, 2.57, 1.72, title="晴天 614.7 kWh", value="88.3%利用", detail="バス163.8 / BESS378.9 / 抑制71.9 kWh", accent=SUN, fill=SUN_LIGHT)
    _add_card(slide, 10.28, 3.16, 2.57, 1.72, title="雨天 101.1 kWh", value="100%利用", detail="バス1.0 / BESS100.1 / 抑制0.0 kWh", accent=RAIN, fill=RAIN_LIGHT)
    _add_text(slide, "PV残差：最大2.7×10⁻¹⁵ kWh", 10.38, 5.25, 2.35, 0.38, size=12.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "晴天PVはEV全車投入ではなく、主にBESS充電と一部直接充電へ配分され、71.9 kWhを抑制。")
    _add_notes(
        slide,
        "目標45秒\nPVの行き先を発電量と重ねました。晴天は614.7キロワット時のうち、163.8を直接バスへ、378.9をBESSへ、71.9を出力抑制しています。雨天は101.1のほぼ全量を使い、主にBESSへ入れています。各時間枠で発電量イコール三つの行き先となり、最大残差は10のマイナス15乗オーダーです。晴天で余るPVがある一方、EV全車使用にならないことも確認できます。",
    )

    # 13. Bus sources
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "EVバス充電の電力源", 13)
    _add_picture(slide, charts["sources"], 0.50, 0.88, 9.55, 5.35)
    _add_card(slide, 10.28, 1.10, 2.57, 1.72, title="晴天 充電入力", value=f"{sunny['daily_energy']['bus_charging_input_kwh']:.1f} kWh", detail="系統1015.6 / PV163.8 / BESS342.0", accent=SUN, fill=SUN_LIGHT)
    _add_card(slide, 10.28, 3.16, 2.57, 1.72, title="雨天 充電入力", value=f"{rain['daily_energy']['bus_charging_input_kwh']:.1f} kWh", detail="系統1000.7 / PV1.0 / BESS90.3", accent=RAIN, fill=RAIN_LIGHT)
    _add_text(slide, "電源合計−充電入力残差\n最大3.41×10⁻¹² kWh", 10.35, 5.18, 2.42, 0.55, size=12.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "PV減少とEV担当便減少が同時に起きるため、雨天の系統購入日量は単純には増えない。")
    _add_notes(
        slide,
        "目標50秒\nEVバス充電の電力源です。晴天の充電入力は1521.4キロワット時で、系統1015.6、PV直接163.8、BESS342.0です。雨天は総充電入力が1092.0まで減り、系統1000.7、PV直接1.0、BESS90.3です。雨天でPVが減っても系統購入日量が大きく増えないのは、EV担当便が22便減って充電需要そのものも小さくなるためです。電源合計と充電入力の差は数値誤差内です。",
    )

    # 14. Grid import
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "系統購入電力とピーク", 14)
    _add_picture(slide, charts["grid"], 0.55, 0.95, 9.45, 4.80)
    _add_card(slide, 10.28, 1.12, 2.57, 1.50, title="晴天", value=f"{sunny['daily_energy']['grid_import_kwh']:.1f} kWh", detail=f"ピーク {sunny['daily_energy']['peak_grid_import_kw']:.1f} kW", accent=SUN, fill=SUN_LIGHT)
    _add_card(slide, 10.28, 2.92, 2.57, 1.50, title="雨天", value=f"{rain['daily_energy']['grid_import_kwh']:.1f} kWh", detail=f"ピーク {rain['daily_energy']['peak_grid_import_kw']:.1f} kW", accent=RAIN, fill=RAIN_LIGHT)
    _add_text(slide, "雨天−晴天\n日量 −14.9 kWh\nピーク +24.8 kW", 10.42, 4.76, 2.28, 0.82, size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "需要料金は日量ではなく最大kWで決まるため、雨天は系統日量が少なくても需要料金が高い。")
    _add_notes(
        slide,
        "目標45秒\n系統購入を60分平均電力で示しています。晴天の日量は1015.6キロワット時、雨天は1000.7で、雨天の方が14.9少ないです。一方、ピークは晴天95.4キロワット、雨天120.2で、雨天が24.8高くなります。需要料金はエネルギー総量ではなく最大電力で決まるため、雨天の需要料金が約992円高いことと整合します。",
    )

    # 15. Energy balance diagram
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "電力需給帳尻の監査", 15)
    _add_flow_column(slide, sunny, 0.48, SUN, SUN_LIGHT)
    _add_flow_column(slide, rain, 6.90, RAIN, RAIN_LIGHT)
    _add_bottom_line(slide, "PV式・充電源式・BESS遷移式は、両日とも1.0×10⁻⁶ kWh以下で成立。", fill=GREEN_LIGHT, color=GREEN)
    _add_notes(
        slide,
        "目標55秒\nここは帳尻を一枚にまとめた図です。上段でPV発電を直接バス、BESS、抑制へ分け、系統購入を示しています。中段でBESSからの供給と、バス充電入力の三電源合計を示します。BESSは初期300から終端300へ戻り、効率損失も別に数えています。PV式、充電源式、BESS遷移式の最大残差はいずれも許容値10のマイナス6乗キロワット時より十分小さく、数理上の需給帳尻は成立しています。",
    )

    # 16. Fuel consistency
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "エンジンバス燃料と運行の整合", 16)
    _add_picture(slide, charts["fuel"], 0.50, 0.92, 8.85, 4.95)
    _add_rect(slide, 9.62, 1.12, 3.22, 4.55, fill=GRAY_LIGHT, line=GRAY, radius=True)
    _add_text(slide, "距離ベース再計算", 9.84, 1.35, 2.78, 0.35, size=17, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"晴天\n{sunny['fuel']['service_distance_km']:.1f} + {sunny['fuel']['intertrip_deadhead_distance_km']:.1f} km\n→ {sunny['fuel']['total_fuel_l']:.3f} L\n→ {sunny['fuel']['reported_fuel_cost_final_jpy']:,.0f}円", 9.88, 1.92, 2.70, 1.30, size=13.5, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"雨天\n{rain['fuel']['service_distance_km']:.1f} + {rain['fuel']['intertrip_deadhead_distance_km']:.1f} km\n→ {rain['fuel']['total_fuel_l']:.3f} L\n→ {rain['fuel']['reported_fuel_cost_final_jpy']:,.0f}円", 9.88, 3.35, 2.70, 1.30, size=13.5, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "燃料単価150円/L\n再計算残差 < 2×10⁻¹⁰円", 9.90, 4.87, 2.66, 0.58, size=12.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "運行距離とは整合するが、給油イベント0件のため“実現給油量”ではなく距離ベース仮勘定。", fill=RED_LIGHT, color=RED)
    _add_notes(
        slide,
        "目標55秒\nエンジンバスの担当便から営業距離と便間回送距離を再集計しました。晴天は合計1287.2キロメートルから284.773リットル、雨天は1538.4キロメートルから340.364リットルとなり、150円毎リットルを掛けた燃料費は報告値と実質ゼロ円差です。ただし、給油イベントは両日ゼロです。したがって、運行距離との整合は確認できましたが、燃料タンク残量と給油タイミングの実現可能性までは今回の結果から主張できません。",
    )

    # 17. Cost comparison
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "費用比較：晴天と雨天", 17)
    _add_picture(slide, charts["cost"], 0.55, 0.92, 9.35, 4.95)
    total_delta = rain["costs_jpy"]["total_cost"] - sunny["costs_jpy"]["total_cost"]
    fuel_delta = rain["costs_jpy"]["fuel_cost"] - sunny["costs_jpy"]["fuel_cost"]
    demand_delta = rain["costs_jpy"]["demand_cost"] - sunny["costs_jpy"]["demand_cost"]
    _add_card(slide, 10.18, 1.12, 2.67, 1.58, title="総費用差", value=f"+{total_delta:,.0f}円", detail=f"雨天は晴天比 +{total_delta/sunny['costs_jpy']['total_cost']*100:.2f}%", accent=RED, fill=RED_LIGHT)
    _add_text(slide, f"主な増加要因\n燃料 +{fuel_delta:,.0f}円\n需要料金 +{demand_delta:,.0f}円", 10.38, 3.05, 2.27, 1.00, size=15, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "電力費はほぼ同額\nEV便数の差とピーク時刻を同時に見る必要", 10.38, 4.47, 2.27, 0.75, size=12.5, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bottom_line(slide, "雨天コスト増の中心は、EV便減少による燃料費と高い系統ピークによる需要料金。")
    _add_notes(
        slide,
        "目標45秒\n会計総額は晴天71万3032円、雨天72万2511円で、雨天が9479円、1.33パーセント高いです。車両使用費64万円が共通で大きいため、右側で変動要因を拡大しています。雨天増加の主因は燃料費8339円と需要料金992円です。電力費がほぼ同額なのは、PV減少と同時にEV担当便も減って系統購入日量が減るためです。",
    )

    # 18. Summary and next actions
    slide = prs.slides.add_slide(blank)
    _add_header(slide, "まとめと次の検証", 18)
    _add_rect(slide, 0.48, 1.00, 6.05, 4.98, fill=BLUE_LIGHT, line=BLUE, radius=True)
    _add_text(slide, "今回答えられること", 0.76, 1.28, 5.48, 0.40, size=20, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(
        slide,
        [
            "BESSは両日とも日末ΔSOC=0 kWh",
            "晴天でもEVは16/35台で、全車使用ではない",
            "PV・充電源・BESSの需給式は許容誤差内",
            "燃料費は営業＋便間回送距離と整合",
            "晴天はEV便が22便多く、総費用が9,479円低い",
        ],
        0.82,
        1.92,
        5.36,
        3.35,
        size=15.2,
    )
    _add_rect(slide, 6.82, 1.00, 6.03, 4.98, fill=SUN_LIGHT, line=SUN, radius=True)
    _add_text(slide, "今後の検証課題", 7.10, 1.28, 5.47, 0.40, size=20, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_bullets(
        slide,
        [
            "ICE26台が正ならscenarioを修正し再計算",
            "燃料タンクSOC・給油時間・給油設備を実現台帳化",
            "変更内容を整理した版で両日を同条件再実行",
            "Stage 1 gapを縮め、複数seedで頑健性を確認",
            "EV全車使用／EV台数制約の感度分析を別実験化",
        ],
        7.16,
        1.92,
        5.34,
        3.35,
        size=15.2,
    )
    _add_bottom_line(slide, "現時点の成果は“内部整合したPhase 3可行解”。正式KPI化には入力台数と最適性の残課題を解消する。")
    _add_notes(
        slide,
        "目標55秒\nまとめです。BESSの日末SOC、EV使用台数、電力需給、燃料距離整合という先生の確認事項には、現結果から定量的に答えられます。一方で、現入力はICE25台であり26台ではありません。また燃料は給油計画ではなく距離仮勘定、Stage 1は約13パーセントのギャップが残る暫定可行解です。次はICE台数の正本確認、燃料タンクと給油の実現台帳化、clean commit後の同条件再実行、そしてEV全車使用を含む感度分析を行います。",
    )

    if len(prs.slides) != TOTAL_SLIDES:
        raise AssertionError(f"Expected {TOTAL_SLIDES} slides, got {len(prs.slides)}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def run(args: argparse.Namespace) -> int:
    audit = _load_json(Path(args.audit).resolve())
    output_path = Path(args.output).resolve()
    build_deck(audit, output_path, Path(args.asset_dir).resolve())
    print(output_path)
    return 0


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--audit", default=str(DEFAULT_AUDIT))
    parser.add_argument("--output", default=str(DEFAULT_PRESENTATION))
    parser.add_argument("--asset-dir", default=str(DEFAULT_ASSET_DIR))
    return run(parser.parse_args())


if __name__ == "__main__":
    raise SystemExit(main())
